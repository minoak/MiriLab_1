# -*- coding: utf-8 -*-
"""압축본 pptx 구조 검증(일회용) — 장수·그림·노트·슬라이드 경계 초과 확인."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent
p = Presentation(ROOT / "notebooks" / "미리랩_발표_요약.pptx")
SW, SH = 13.333, 7.5
print("slides:", len(p.slides))
for i, s in enumerate(p.slides, 1):
    n_pic = sum(1 for sh in s.shapes if sh.shape_type == 13)
    has_note = (bool(s.notes_slide.notes_text_frame.text.strip())
                if s.has_notes_slide else False)
    overflow = []
    for sh in s.shapes:
        try:
            r = Emu(sh.left + sh.width).inches
            b = Emu(sh.top + sh.height).inches
            if r > SW + 0.01 or b > SH + 0.01:
                overflow.append(sh.shape_id)
        except TypeError:
            pass
    print(i, "shapes=%d" % len(s.shapes), "pics=%d" % n_pic,
          "note=%s" % has_note,
          ("OVERFLOW:" + str(overflow)) if overflow else "bounds-ok")
