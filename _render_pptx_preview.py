# -*- coding: utf-8 -*-
"""pptx 근사 미리보기 렌더러 (PowerPoint 없이 레이아웃 검수용).

사용법: python _render_pptx_preview.py
출력:   notebooks/_pptx_assets/preview_01.png ... preview_19.png

PowerPoint 의 실제 렌더링과 100% 같지는 않지만(폰트 메트릭 근사),
텍스트 넘침·도형 겹침·정렬 문제를 잡기에 충분하다.
텍스트가 박스를 넘치면 박스 둘레에 빨간 점선을 그려 표시한다.
"""
from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent
PPTX = ROOT / "notebooks" / "미리랩_발표.pptx"
OUTDIR = ROOT / "notebooks" / "_pptx_assets"

DPI = 110  # px per inch
EMU_PER_IN = 914400

FONT_REG = r"C:\Windows\Fonts\malgun.ttf"
FONT_BOLD = r"C:\Windows\Fonts\malgunbd.ttf"
FONT_EMOJI = r"C:\Windows\Fonts\seguiemj.ttf"

_font_cache: dict = {}


def get_font(size_px: int, bold: bool, emoji: bool = False):
    key = (size_px, bold, emoji)
    if key not in _font_cache:
        if emoji:
            try:
                _font_cache[key] = ImageFont.truetype(FONT_EMOJI, size_px)
            except Exception:
                _font_cache[key] = ImageFont.truetype(FONT_REG, size_px)
        else:
            _font_cache[key] = ImageFont.truetype(FONT_BOLD if bold else FONT_REG,
                                                  size_px)
    return _font_cache[key]


def emu2px(v) -> int:
    return int(round(v / EMU_PER_IN * DPI))


def is_emoji(ch: str) -> bool:
    cp = ord(ch)
    return (0x1F000 <= cp <= 0x1FAFF) or (0x2600 <= cp <= 0x27BF) or cp in (
        0x2B05, 0x2B06, 0x2B07, 0x2B50, 0x20E3, 0xFE0F, 0x23F8, 0x25B6)


def tokenize(text: str):
    """ASCII 단어는 묶고 CJK/기타는 글자 단위로."""
    tokens, cur = [], ""
    for ch in text:
        if ch.isascii() and not ch.isspace():
            cur += ch
        else:
            if cur:
                tokens.append(cur)
                cur = ""
            tokens.append(ch)
    if cur:
        tokens.append(cur)
    return tokens


def run_style(run):
    f = run.font
    size_pt = f.size.pt if f.size is not None else 18
    bold = bool(f.bold)
    color = (31, 41, 55)
    try:
        if f.color and f.color.rgb is not None:
            c = f.color.rgb
            color = (c[0], c[1], c[2]) if isinstance(c, (tuple, list)) else (
                (int(str(c)[0:2], 16), int(str(c)[2:4], 16), int(str(c)[4:6], 16)))
    except Exception:
        pass
    return size_pt, bold, color


def measure(draw, text, font):
    try:
        return draw.textlength(text, font=font)
    except Exception:
        return font.getbbox(text)[2]


def layout_paragraph(draw, runs, max_w):
    """runs: [(text, font, color, size_px)] -> lines: [[(seg, font, color)], ...], line_heights"""
    lines, cur, cur_w = [], [], 0.0
    max_size = 1
    line_sizes = []

    def flush():
        nonlocal cur, cur_w, max_size
        lines.append(cur)
        line_sizes.append(max_size if cur else 1)
        cur, cur_w, max_size = [], 0.0, 1

    for text, font, color, size_px in runs:
        for tok in tokenize(text):
            w = measure(draw, tok, font)
            if cur and cur_w + w > max_w and tok.strip():
                flush()
                if tok == " ":
                    continue
            if w > max_w and len(tok) > 1:  # 한 토큰이 너무 길면 글자 분해
                for ch in tok:
                    cw = measure(draw, ch, font)
                    if cur and cur_w + cw > max_w:
                        flush()
                    cur.append((ch, font, color))
                    cur_w += cw
                    max_size = max(max_size, size_px)
            else:
                cur.append((tok, font, color))
                cur_w += w
                max_size = max(max_size, size_px)
    if cur or not lines:
        flush()
    return lines, line_sizes


def render_text_frame(draw, img, shape, ox, oy, w, h):
    tf = shape.text_frame
    ml = emu2px(tf.margin_left) if tf.margin_left is not None else 8
    mr = emu2px(tf.margin_right) if tf.margin_right is not None else 8
    mt = emu2px(tf.margin_top) if tf.margin_top is not None else 4
    mb = emu2px(tf.margin_bottom) if tf.margin_bottom is not None else 4
    inner_w = max(10, w - ml - mr)

    # 전체 문단 레이아웃 먼저 계산 (anchor 처리용)
    blocks = []  # (lines, line_sizes, align, line_spacing, space_after, space_before)
    total_h = 0
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            size_pt, bold, color = run_style(r)
            size_px = max(6, int(round(size_pt * DPI / 72)))
            # 이모지 구간 분리
            buf = ""
            for ch in r.text:
                if is_emoji(ch):
                    if buf:
                        runs.append((buf, get_font(size_px, bold), color, size_px))
                        buf = ""
                    runs.append((ch, get_font(size_px, bold, emoji=True),
                                 color, size_px))
                else:
                    buf += ch
            if buf:
                runs.append((buf, get_font(size_px, bold), color, size_px))
        if not runs:
            blocks.append(([], [], PP_ALIGN.LEFT, 1.0, 4, 0))
            total_h += 8
            continue
        lines, line_sizes = layout_paragraph(draw, runs, inner_w)
        ls = p.line_spacing if p.line_spacing else 1.0
        sa = p.space_after.pt if p.space_after is not None else 0
        sb = p.space_before.pt if p.space_before is not None else 0
        align = p.alignment if p.alignment is not None else PP_ALIGN.LEFT
        blocks.append((lines, line_sizes, align, ls, sa, sb))
        for s in line_sizes:
            total_h += int(s * 1.25 * ls)
        total_h += int((sa + sb) * DPI / 72)

    anchor = tf.vertical_anchor
    y = oy + mt
    avail_h = h - mt - mb
    if anchor == MSO_ANCHOR.MIDDLE and total_h < avail_h:
        y = oy + mt + (avail_h - total_h) // 2
    overflow = total_h > avail_h + 6

    for lines, line_sizes, align, ls, sa, sb in blocks:
        y += int(sb * DPI / 72)
        if not lines:
            y += 8
            continue
        for line, lsize in zip(lines, line_sizes):
            line_w = sum(measure(draw, seg, f) for seg, f, _ in line)
            if align == PP_ALIGN.CENTER:
                x = ox + ml + (inner_w - line_w) / 2
            elif align == PP_ALIGN.RIGHT:
                x = ox + ml + inner_w - line_w
            else:
                x = ox + ml
            line_h = int(lsize * 1.25 * ls)
            for seg, font, color in line:
                try:
                    draw.text((x, y), seg, font=font, fill=color,
                              embedded_color=True)
                except Exception:
                    draw.text((x, y), seg, font=font, fill=color)
                x += measure(draw, seg, font)
            y += line_h
        y += int(sa * DPI / 72)

    if overflow:
        for i in range(0, 2 * (w + h), 14):  # 빨간 점선 표시
            pass
        draw.rectangle([ox, oy, ox + w, oy + h], outline=(255, 0, 0), width=3)
        draw.text((ox + 2, oy + 2), "OVERFLOW", font=get_font(14, True),
                  fill=(255, 0, 0))


def shape_colors(shape):
    fill_rgb = line_rgb = None
    try:
        if shape.fill.type is not None and str(shape.fill.type) != "MSO_FILL_TYPE.BACKGROUND (5)":
            c = shape.fill.fore_color.rgb
            if c is not None:
                s = str(c)
                fill_rgb = (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        pass
    try:
        lf = shape.line.fill
        if lf.type is not None and "BACKGROUND" not in str(lf.type):
            c = shape.line.color.rgb
            if c is not None:
                s = str(c)
                line_rgb = (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        pass
    return fill_rgb, line_rgb


def render_slide(slide, sw_px, sh_px):
    img = Image.new("RGB", (sw_px, sh_px), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        try:
            x, y = emu2px(shape.left), emu2px(shape.top)
            w, h = emu2px(shape.width), emu2px(shape.height)
        except Exception:
            continue
        st = str(shape.shape_type)
        if "PICTURE" in st:
            try:
                pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                pic = pic.resize((max(1, w), max(1, h)))
                img.paste(pic, (x, y), pic)
            except Exception:
                draw.rectangle([x, y, x + w, y + h], outline=(200, 0, 200),
                               width=2)
            continue
        if "AUTO_SHAPE" in st or "ROUND" in st or "OVAL" in st or "RECTANGLE" in st:
            fill_rgb, line_rgb = shape_colors(shape)
            if "OVAL" in str(getattr(shape, "auto_shape_type", "")):
                if fill_rgb:
                    draw.ellipse([x, y, x + w, y + h], fill=fill_rgb,
                                 outline=line_rgb)
                elif line_rgb:
                    draw.ellipse([x, y, x + w, y + h], outline=line_rgb)
            else:
                radius = 10 if "ROUNDED" in str(getattr(shape, "auto_shape_type", "")) else 0
                if fill_rgb or line_rgb:
                    draw.rounded_rectangle([x, y, x + w, y + h], radius=radius,
                                           fill=fill_rgb, outline=line_rgb,
                                           width=2 if line_rgb else 0)
            if shape.has_text_frame and shape.text_frame.text.strip():
                render_text_frame(draw, img, shape, x, y, w, h)
            continue
        if shape.has_text_frame:
            render_text_frame(draw, img, shape, x, y, w, h)
    return img


def main():
    prs = Presentation(PPTX)
    sw_px = emu2px(prs.slide_width)
    sh_px = emu2px(prs.slide_height)
    OUTDIR.mkdir(exist_ok=True)
    for i, slide in enumerate(prs.slides, 1):
        img = render_slide(slide, sw_px, sh_px)
        p = OUTDIR / f"preview_{i:02d}.png"
        img.save(p)
        print(f"[OK] {p.name}")
    print("[DONE]", len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    main()
