# -*- coding: utf-8 -*-
"""미리랩 발표자료 — 압축본(.pptx) 생성 스크립트.

사용법:  python _make_pptx_short.py
출력:    notebooks/미리랩_발표_요약.pptx  (14슬라이드, 발표자 노트 포함)

풀버전(`_make_pptx.py` → 미리랩_발표.pptx, 23장)과의 관계:
- 풀버전 = 자료집(제출·아카이브용, 모든 시행착오·검증 상세). **그대로 유지.**
- 이 파일 = 발표 시나리오용 압축본. 한 장 = 한 메시지, 장당 35~45초
  + 미리마을 데모 60~90초  →  총 ~10분.
- 컷 기준: 화면(데모)이 보이는 장은 살리고, 글로 설명하는 장은 합치거나 자료집으로.
  (아젠다·협업규율 삭제 / 문제+컨셉 합본 / 시스템 3장→1장 / 시행착오 4장→1장 /
   검증 3장→2장 / 확장 3장→2장 / 한계 5→3)
- 풀버전에 없는 신설 1장: 활용 방안(12번 — 정책 수명주기 프레임, 실제 탭 기능 근거).

헬퍼·팔레트·타이틀·클로징은 `_make_pptx`에서 그대로 가져온다(단일 소스).
스크린샷 자리는 풀버전과 같은 회색 점선 박스(📸) — 캡처 후 교체.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from _make_pptx import (
    BLUE, CARD, GRAY_BG, GRAY_BORDER, GREEN, INK, LINE, NAVY, ORANGE,
    PALE_GREEN, PALE_NAVY, PALE_ORANGE, PALE_RED, RED, ROOT, SUB, SW, SH,
    arrow_text, box, bullets, flow_box, header, img_fit, note, para,
    prep_sprites, role_card, s01_title, s19_closing, shot, tb,
)

OUT = ROOT / "notebooks" / "미리랩_발표_요약.pptx"


# ================================================================ 슬라이드들
def t02_problem_concept(prs):
    """문제 + 컨셉 합본 — 풀버전 3·4장을 한 장으로."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 2, "01 문제", "정책은 왜 닿지 못하는가 — 같은 정책, 다른 인생",
           title_size=24)
    # 사각지대 4유형 — 한 줄 미니카드
    cards = [
        ("정보 격차", "대상자인지조차 모른다"),
        ("이해도 부족", "문구·신청 조건이 어렵다"),
        ("디지털 장벽", "온라인 신청 자체가 벽"),
        ("도달 실패", "발표돼도 닿지 않는다"),
    ]
    x = 0.6
    for title, body in cards:
        b = box(slide, x, 1.7, 2.93, 0.95, fill=PALE_RED, line=RED, lw=1.0)
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, title, size=13.5, bold=True, color=RED, first=True, after=2,
             align=PP_ALIGN.CENTER)
        para(tf, body, size=11, color=INK, after=0, align=PP_ALIGN.CENTER)
        x += 3.07
    # 빈자리 한 줄
    box(slide, 0.6, 2.85, 12.13, 0.95, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tfb = tb(slide, 0.95, 2.97, 11.5, 0.72, anchor=MSO_ANCHOR.MIDDLE)
    para(tfb, [("기존 도구(정부24 챗봇·복지로)는 정책을 ", {}),
               ("‘찾는’ 시민", {"bold": True, "color": BLUE}),
               ("을 돕는다 — 정책을 ", {}),
               ("‘만드는’ 입안자", {"bold": True, "color": RED}),
               ("가 “누가 못 받을까?”를 배포 전에 볼 도구는 없다",
                {"bold": True})], size=14, first=True, after=0,
         align=PP_ALIGN.CENTER)
    # 핵심 장면 — 대조 3명
    t = tb(slide, 0.6, 4.0, 12.13, 0.4)
    para(t, "핵심 장면 — 같은 ‘청년 월세 지원’ 정책, 6개월 후", size=14,
         bold=True, color=INK, first=True, after=0)
    role_card(slide, 0.6, 4.45, 3.93, 1.6, GREEN, PALE_GREEN,
              "수혜  ·  A (26세 청년)",
              ["복지로에서 5분 만에 신청,", "첫 저축을 시작했다."], body_size=12)
    role_card(slide, 4.7, 4.45, 3.93, 1.6, ORANGE, PALE_ORANGE,
              "경계  ·  B (46세 중장년)",
              ["본인은 대상이 아니지만", "자녀를 위해 대신 신청해 준다."],
              body_size=12)
    role_card(slide, 8.8, 4.45, 3.93, 1.6, RED, PALE_RED,
              "사각지대  ·  C (74세 고령)",
              ["6개월이 지나도", "존재조차 알지 못했다."], body_size=12)
    t2 = tb(slide, 0.6, 6.3, 12.13, 0.55)
    para(t2, "“같은 정책, 다른 인생.” — 미리랩은 이 갈림을 배포 전에 보여준다",
         size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER, first=True,
         after=0)
    note(slide, "정부는 매년 수천 건의 정책을 발표하지만 정작 필요한 사람에게 닿지 "
         "못합니다 — 모르거나, 어렵거나, 디지털이 벽이거나, 발표돼도 닿지 않거나. "
         "기존 서비스는 전부 정책을 '찾는' 시민용이고, 정책을 '만드는' 입안자가 "
         "'누가 못 받을지'를 배포 전에 보는 도구는 없습니다. 저희 답이 미리랩입니다. "
         "같은 청년 월세 정책이라도 26세는 5분 만에 신청하고, 46세는 자녀를 위해 대신 "
         "신청하고, 74세는 존재조차 모릅니다. '같은 정책, 다른 인생' — 이 갈림을 배포 "
         "전에 보여드립니다. 수혜·경계·사각 세 가지 색이 오늘 발표를 관통합니다.")


def t03_system(prs):
    """시스템 한 장 — 풀버전 5·6·7장(파이프라인+데이터셋+페르소나)을 한 장으로."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 3, "02 시스템", "통계에서 뽑은 시민 24명 — 3축 단방향 파이프라인",
           title_size=23)
    # 바닥 재료: 페르소나
    flow_box(slide, 0.6, 1.7, 12.13, 1.15,
             "페르소나 — nvidia/Nemotron-Personas-Korea  (실제 한국 통계 기반 합성 · 100만 명 × 26항목 · CC BY 4.0)",
             ["seed 고정 24명 샘플링 — 손으로 만들면 보고 싶은 반응을 심게 된다. "
              "실제 분포에서 뽑으면 사각지대 인물이 처음부터 표본 안에 들어 있다."],
             fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY, tsize=13, bsize=11.5)
    arrow_text(slide, 6.42, 2.9, 0.5, "↓", size=15)
    # 3축 단방향
    flow_box(slide, 1.0, 3.3, 3.5, 1.5, "축1 · 정보",
             ["react — 시민 각자의 t0 반응", "점수는 익명 설문 — LLM은",
              "선택지만, 숫자 변환은 코드"],
             fill=PALE_GREEN, line_c=GREEN, tcolor=GREEN, tsize=14, bsize=10.8)
    flow_box(slide, 4.92, 3.3, 3.5, 1.5, "축2 · 결과",
             ["인생극장 — 전원 1·3·6개월", "시간경과 시뮬, 실제 결과에서",
              "대조 3명(수혜·경계·사각) 선별"],
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=14,
             bsize=10.8)
    flow_box(slide, 8.83, 3.3, 3.5, 1.5, "축3 · 요약",
             ["집계·서술 (axis3 · 순수 함수)", "게이지·리포트 등 화면의",
              "모든 숫자가 여기 한 곳에서"],
             fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY, tsize=14, bsize=10.8)
    arrow_text(slide, 4.42, 3.85, 0.5, "→")
    arrow_text(slide, 8.33, 3.85, 0.5, "→")
    n = tb(slide, 1.0, 4.95, 11.3, 0.35)
    para(n, "단방향이라 뒤로 흐르지 않고, 판정 지점이 하나라 모순이 생기지 않는다 "
         "(v1 단일 LangGraph의 축 간 모순 → 재설계)", size=11.5, color=SUB,
         first=True, after=0)
    arrow_text(slide, 6.42, 5.32, 0.5, "↓", size=15)
    flow_box(slide, 1.0, 5.7, 11.33, 0.6,
             "Streamlit 7탭  ·  데모 = Gemini 실런 녹화 재생 (발표장 네트워크 0콜)",
             None, fill=CARD, line_c=NAVY, tsize=12.5)
    f = tb(slide, 1.0, 6.5, 11.33, 0.4)
    para(f, "재현성 — seed 고정 + 로컬 캐시: 발표마다 같은 시민 24명, 오프라인에서도 동작",
         size=10.5, color=SUB, first=True, after=0)
    note(slide, "시스템은 한 장이면 충분합니다. 바닥엔 NVIDIA가 실제 한국 통계로 합성해 "
         "공개한 페르소나 100만 명이 있고, 저희는 seed 고정으로 24명을 뽑습니다 — "
         "손으로 만들면 보고 싶은 반응을 심게 되지만, 실제 분포에서 뽑으면 사각지대 "
         "인물이 처음부터 표본 안에 들어 있습니다. 그 위가 3축 단방향 파이프라인입니다. "
         "축1이 시민 각자의 첫 반응을 만들고, 축2가 그걸 시딩 삼아 시간 속에서 인생을 "
         "굴리고, 축3가 두 축을 읽어 집계합니다. 단방향이라 뒤로 흐르지 않고 판정 "
         "지점이 하나라 모순이 안 생깁니다 — 처음 단일 LangGraph 시절의 축 간 모순을 "
         "재설계로 푼 결과입니다. 데모는 Gemini 실런을 녹화해 재생하므로 발표장 "
         "네트워크가 죽어도 진짜 결과로 시연됩니다. 이제 세 축을 화면으로 하나씩 "
         "보여드리겠습니다.")


def t04_axis1(prs):
    """축1 — 풀버전 8장(번호·노트만 압축)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 4, "02 시스템", "축1 · 정보 — 시민 반응", accent=GREEN)
    fx, fw = 0.6, 5.8
    para(tb(slide, fx, 1.75, fw, 0.6),
         "시민 각자가 정책을 보고 즉시 반응한다 (t0)", size=15, bold=True,
         color=NAVY, first=True, after=0)
    flow_box(slide, fx, 2.5, fw, 1.55, "react — 1차 반응",
             ["· 반응문 + 입장 (찬성 · 반대 · 혼합)",
              "· 익명 설문 응답",
              "· LLM은 선택지만 고르고, 0~100 점수 변환은 코드"],
             fill=PALE_GREEN, line_c=GREEN, tcolor=GREEN, tsize=14, bsize=12)
    box(slide, fx, 4.32, fw, 1.05, fill=GRAY_BG, line=GRAY_BORDER, dash=True)
    tfi = tb(slide, fx + 0.22, 4.43, fw - 0.44, 0.85)
    para(tfi, "interact — 다른 시민을 보고 입장 변화", size=12.5, bold=True,
         color=SUB, first=True, after=2)
    para(tfi, "SNS 채팅 탭 데모 전용 · 집계엔 미반영 (메인 라인에서 덜어냄)",
         size=11.5, color=SUB, after=0)
    para(tb(slide, fx, 5.65, fw, 0.5),
         "→ 집단의 첫 온도. 게이지·리포트 집계는 축3에서.", size=12.5,
         bold=True, color=GREEN, first=True, after=0)
    shot(slide, 6.7, 1.75, 6.03, 5.05, "‘시민 반응’ 탭 캡처",
         "데모에서 실제 반응이 채워지는 화면")
    note(slide, "축1은 시민 각자의 첫 반응(t0)입니다. 정책을 보고 반응문과 입장, "
         "익명 설문에 답합니다. 중요한 건 — LLM은 0~100 점수를 직접 매기지 않습니다. "
         "설문 선택지만 고르고 숫자 변환은 코드가 합니다. 구버전엔 점수를 직접 물었는데 "
         "기분이 점수에 번지는 정서 후광 탓에 비대상 노인이 청년 정책에 수혜 70을 주는 "
         "버그가 있었고, 설문 전환으로 해소했습니다. 화면의 게이지·집계는 이 단계가 "
         "아니라 축3에서 나옵니다.")


def t05_axis2(prs):
    """축2 — 풀버전 9장(번호·노트만 압축)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 5, "02 시스템", "축2 · 결과 — 정책 인생극장", accent=ORANGE)
    fx, fw = 0.6, 5.6
    flow_box(slide, fx, 1.8, fw, 0.86, "① 전원 시간경과 시뮬 (1 · 3 · 6개월)",
             "시점마다: 장소 + 경로 + 행동 서사 + 상태 + 막힌 지점",
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=12.5,
             bsize=10.5)
    arrow_text(slide, fx + fw / 2 - 0.6, 2.68, 1.2, "↓", size=14)
    flow_box(slide, fx, 3.02, fw, 0.86, "② 실제 결과에서 대조 3명 선별",
             "받음 → 수혜 / 진행 중 → 경계 / 막힘·못 닿음 → 사각",
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=12.5,
             bsize=10.5)
    arrow_text(slide, fx + fw / 2 - 0.6, 3.9, 1.2, "↓", size=14)
    flow_box(slide, fx, 4.24, fw, 0.95, "③ 접근 여정 카드",
             "🧭 알게됨 → 신청 → ⛔막힘  (경로·막힌 지점·삶의 변화)",
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=12.5,
             bsize=10.5)
    tf = tb(slide, fx, 5.42, fw, 1.4)
    para(tf, [("채널 5곳이 사각지대의 표현 장치  ",
               {"bold": True, "color": NAVY, "size": 13})], first=True, after=4)
    para(tf, "복지로(온라인) · 주민센터 · 복지관 · 직장/시장 · 집", size=12,
         color=INK, after=3)
    para(tf, "청년은 복지로에서 5분, 고령은 어느 채널에도 닿지 못한다.",
         size=12, bold=True, color=RED, after=0)
    shot(slide, 6.5, 1.8, 6.23, 5.0, "‘정책 인생극장’ 탭 캡처",
         "대조 3명 카드 + 여정 띠(알게됨→신청→막힘)가 보이는 화면 추천")
    note(slide, "축2 정책 인생극장입니다. 시민 전원을 1·3·6개월 시간 경과 속에서 살게 "
         "하고, 그 '실제 결과'에서 수혜·경계·사각 세 사람을 뽑아 카드로 보여줍니다. "
         "카드를 펼치면 '알게 됨 → 신청 → 서류에서 막힘' 같은 여정 띠와 함께 어떤 "
         "경로로 알게 됐고 어디서 막혔는지가 서사로 나옵니다. 미리 점수로 3명을 "
         "'예측'해 뽑던 구버전은 라벨과 이야기가 어긋나는 모순이 있어서, 전원을 먼저 "
         "살게 한 뒤 결과에서 뽑는 방식으로 바꿨습니다 — 수혜·경계·사각은 점수가 "
         "아니라 결과니까요. 청년에겐 복지로에서 5분인 일이 고령 어르신에겐 어느 "
         "채널로도 닿지 못하는 일이 됩니다.")


def t06_axis3(prs):
    """축3 — 풀버전 10장(번호·노트만 압축)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 6, "02 시스템", "축3 · 요약 — 집계", accent=NAVY)
    fx, fw = 0.6, 5.8
    para(tb(slide, fx, 1.75, fw, 0.6),
         "axis3가 축1·축2를 다시 읽어 집계한다 (순수 함수 · LLM 무관)",
         size=15, bold=True, color=NAVY, first=True, after=0)
    flow_box(slide, fx, 2.5, fw, 1.35, "게이지 3종",
             ["정책수용도 · 신청의향지수 · 사회혼란도(=불만 평균)"],
             fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY, tsize=14, bsize=12)
    flow_box(slide, fx, 4.0, fw, 1.35, "종합 리포트 (.md)",
             ["숫자·인용은 코드가 채우고, 산문 4칸만 LLM (1콜)"],
             fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY, tsize=14, bsize=12)
    box(slide, fx, 5.55, fw, 0.85, fill=PALE_GREEN, line=GREEN)
    para(tb(slide, fx + 0.22, 5.67, fw - 0.44, 0.65),
         "화면에 뜨는 모든 숫자는 여기서 — 판정 지점이 하나라 모순 없음",
         size=12.5, bold=True, color=GREEN, first=True, after=0)
    shot(slide, 6.7, 1.75, 6.03, 5.05, "게이지 · 리포트 화면 캡처",
         "게이지 3종 + 종합 리포트가 보이는 화면")
    note(slide, "축3는 집계입니다. 핵심은 — 화면에 뜨는 게이지와 리포트의 모든 숫자가 "
         "여기 한 곳에서 나온다는 점입니다. 축1의 반응과 축2의 인생극장 결과를 입력으로 "
         "받아 순수 파이썬 함수로 집계하므로, 같은 입력이면 늘 같은 숫자가 나옵니다. "
         "종합 리포트는 숫자·인용을 코드가 채우고 LLM은 산문 네 칸만 쓰는 고정 "
         "양식입니다. 집계를 한 곳으로 모은 게 '판정 지점이 하나라 모순이 안 생긴다'는 "
         "재설계의 핵심입니다.")


def t07_lessons(prs):
    """시행착오 — 풀버전 11~14장을 대표 2가지 + 한 줄로."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 7, "03 시행착오", "막다른 길이 설계 원칙이 됐다 — 대표 2가지",
           accent=ORANGE, title_size=24)
    # 좌: 10억 버그 → 프롬프트로 조향하지 마라
    box(slide, 0.6, 1.75, 5.9, 4.25, fill=PALE_ORANGE, line=ORANGE, lw=1.0)
    tf1 = tb(slide, 0.85, 1.95, 5.4, 3.9)
    para(tf1, "① “전 국민 10억 지급” 버그", size=15, bold=True, color=ORANGE,
         first=True, after=7)
    bullets(tf1, [
        ("극단 정책에 전원 찬성 — 시민이 ‘이해득실 계산기’", {}),
        ("처방 “현실적으로 판단하라” → 전원 반대로 해결된 듯", {}),
        ([("그러나 정상 정책의 혼란도까지 오르는 ", {}),
          ("과교정", {"bold": True, "color": RED}),
          (" → 판단 지시 전부 철거 (프롬프트는 무대만)", {})], {}),
        ("진짜 답은 지시가 아니라 측정(익명 설문)과 모델 선택(Gemini)", {}),
    ], size=12, after=7, first_done=True, lh=1.25)
    para(tf1, "", size=4, after=2)
    para(tf1, [("교훈 — 프롬프트로 조향하지 마라. ", {"bold": True,
                                                  "color": NAVY}),
               ("분포는 데이터와 측정으로 푼다.", {"color": INK})], size=13,
         after=0, lh=1.25)
    # 우: 판단은 LLM에게, 불변식은 코드로
    box(slide, 6.83, 1.75, 5.9, 4.25, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tf2 = tb(slide, 7.08, 1.95, 5.4, 3.9)
    para(tf2, "② 판단은 LLM에게, 불변식은 코드로", size=15, bold=True,
         color=NAVY, first=True, after=7)
    para(tf2, "LLM은 확률적 — 가끔 말이 안 되는 출력을 낸다. 반복해서 만난 답:",
         size=12, color=INK, after=7, lh=1.25)
    bullets(tf2, [
        ([("상태 비퇴행 가드", {"bold": True}),
          ("  ‘신청→모름’ 퇴행은 ‘막힘’으로 자동 교정", {})], {}),
        ([("태그 대상 게이트", {"bold": True}),
          ("  대상 여부 = 태그가 정하는 사실, 결과 = LLM 궤적", {})], {}),
        ([("전파 게이트", {"bold": True}),
          ("  모르던 사람은 ‘아는 사람을 만났을 때만’ 안다", {})], {}),
    ], size=12, after=7, first_done=True, lh=1.25)
    para(tf2, "", size=4, after=2)
    para(tf2, [("서사와 판단은 LLM이, ", {"color": INK}),
               ("어겨선 안 되는 규칙은 코드 가드가 보증.", {"bold": True,
                                                     "color": NAVY})],
         size=13, after=0, lh=1.25)
    # 하단: 나머지 시행착오 한 줄
    box(slide, 0.6, 6.25, 12.13, 0.72, fill=CARD, line=LINE)
    tfb = tb(slide, 0.88, 6.33, 11.6, 0.56, anchor=MSO_ANCHOR.MIDDLE)
    para(tfb, [("그 외 — ", {"bold": True, "color": SUB}),
               ("‘두 심판 문제’(라벨과 서사의 모순 → 결과 기반 선별로 통일) · "
                "8기능 구상 → 검증 가능한 핵심으로 스코프 축소", {"color": INK}),
               ("   (상세: 자료집 11~14장)", {"color": SUB})], size=11.5,
         first=True, after=0)
    note(slide, "이 모양에 오기까지의 시행착오 중 두 가지만 소개합니다. 첫째, 극단 정책 "
         "'전 국민 10억 지급'에 전원 찬성이 나왔습니다. '현실적으로 판단하라'를 넣자 "
         "해결된 듯 보였지만 정상 정책의 혼란도까지 끌어올리는 과교정이 생겨, 판단 "
         "지시를 전부 철거했습니다. 진짜 답은 지시가 아니라 측정 방식과 모델 선택이었고 "
         "— '프롬프트로 조향하지 마라'가 최종 교훈입니다. 둘째, LLM은 확률적이라 "
         "신청했던 사람이 '몰랐다'로 퇴행하는 식의 출력을 냅니다. 저희 답은 일관됩니다 "
         "— 서사와 판단은 LLM에게 맡기되, 어겨선 안 되는 규칙은 코드 가드가 보증한다. "
         "이 두 원칙이 시스템 전체를 관통하고, 다음 장의 검증이 이 위에서 돌아갑니다.")


def t08_validation(prs):
    """검증 개요 + ① 재료 — 풀버전 15장(번호·노트만 압축)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 8, "04 검증", "‘맞혔나’가 아니라 ‘합리적으로 작동하는가’",
           accent=GREEN, title_size=24)
    box(slide, 0.6, 1.62, 12.13, 1.16, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tfb = tb(slide, 0.95, 1.72, 11.5, 0.98, anchor=MSO_ANCHOR.MIDDLE)
    para(tfb, [("미리랩은 미래를 맞히는 예측기가 아니라, 정책의 파장을 미리 보는 시뮬레이터다. ",
                {"bold": True, "color": NAVY, "size": 13}),
               ("‘여론조사를 얼마나 맞혔나’로 평가하지 않는다 — 그 지표는 약하고, 솔직히 인정한다.",
                {"size": 12.5, "color": INK})], first=True, after=3, lh=1.25)
    para(tfb, "대신 — 극단 정책에도 시민이 상식적으로 반응하면, 일반 정책의 반응도 믿을 수 "
         "있다. 기준은 실행 전에 등록하고, 결과가 나빠도 고치지 않는다(사전등록).",
         size=12.5, bold=True, color=GREEN, after=0, lh=1.25)
    rows = [
        ("① 재료", "시민 24명이 한국 인구를 닮았나?", "분포 7/7 통과 — 무작위 추출 "
         "기대범위 안 (부트스트랩 1만회 · LLM 0콜)", GREEN, PALE_GREEN),
        ("② 작동", "반응 차이를 만드는 게 페르소나인가?", "사전등록 0/4 → 부검 반전 "
         "(다음 장)", ORANGE, PALE_ORANGE),
        ("③ 강건성", "극단 정책에도 상식을 지키나?", "행동 벤치마크 7/7 통과 "
         "(다음 장)", GREEN, PALE_GREEN),
    ]
    x = 0.6
    for tag, q, res, c, pale in rows:
        box(slide, x, 2.92, 3.93, 1.5, fill=pale, line=c, lw=1.0)
        tf = tb(slide, x + 0.2, 3.06, 3.55, 1.25)
        para(tf, [(tag + "  ", {"bold": True, "color": c, "size": 13.5}),
                  (q, {"size": 11.5, "color": INK})], first=True, after=4,
             lh=1.2)
        para(tf, res, size=11, color=INK, after=0, lh=1.22)
        x += 4.1
    img_fit(slide, ROOT / "eval" / "persona_eval_viz.png", 0.6, 4.5, 12.13, 2.45)
    note(slide, "'이걸 믿을 수 있나'에 답합니다. 먼저 분명히 — 미리랩은 미래를 맞히는 "
         "예측기가 아니라 정책의 파장을 미리 보는 시뮬레이터고, '여론조사를 맞혔나' "
         "같은 지표는 약할 수밖에 없음을 솔직히 인정합니다. 대신 기댄 논리는 스트레스 "
         "테스트입니다 — 말도 안 되는 극단 정책에도 시민이 상식적으로 반응한다면, 평범한 "
         "정책의 반응도 믿을 수 있다. 그 위에서 재료·작동·강건성 세 가지를, 판정 기준을 "
         "실행 전에 등록하는 사전등록 방식으로 검증했습니다. 아래 그림이 ① 재료입니다. "
         "전체 100만 명 분포(회색)와 우리 24명(파랑)을 변수별로 비교해 부트스트랩 1만 "
         "번으로 판정했고, 7개 변수 전부 통과 — 쏠림 없이 뽑았습니다. ②와 ③은 다음 "
         "장에서 이야기로 보여드립니다.")


def t08b_audit(prs):
    """검증 ②③ — 풀버전 16·17장(부검 반전 + 행동 벤치마크)을 한 장으로."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 9, "04 검증", "사전등록 0/4의 반전 · 행동 벤치마크 7/7",
           accent=GREEN, title_size=24)
    # 좌상: 부검 스토리
    box(slide, 0.6, 1.7, 5.9, 2.75, fill=PALE_ORANGE, line=ORANGE, lw=1.0)
    tf1 = tb(slide, 0.85, 1.86, 5.4, 2.45)
    para(tf1, "② 작동 — 0/4 탈락, 부검해 보니", size=14, bold=True,
         color=ORANGE, first=True, after=6)
    bullets(tf1, [
        ("실험 — 같은 정책을 인물 카드 ON/OFF(익명)로 비교", {}),
        ("결과 — 사전등록 0/4 탈락 · 점수는 그대로 기록", {}),
        ("부검 — 시민들이 원문의 숨은 요건(별도거주·무주택·소득)까지 읽고 "
         "전원 ‘나는 비대상’ 정확 판정", {}),
    ], size=11.5, after=5, first_done=True, lh=1.22)
    para(tf1, [("시뮬이 채점표보다 현실적이었다 ", {"bold": True, "color": NAVY}),
               ("(실제 청년월세도 신청자 2/3가 요건 탈락)", {"color": SUB})],
         size=12, after=0, lh=1.2)
    # 좌하: 시민 인용
    box(slide, 0.6, 4.6, 5.9, 1.8, fill=CARD, line=LINE)
    tfq = tb(slide, 0.85, 4.74, 5.4, 1.55)
    para(tfq, "시민들이 직접 말한 탈락 사유", size=12.5, bold=True, color=NAVY,
         first=True, after=5)
    para(tfq, [("심석현(26)  ", {"bold": True, "color": BLUE, "size": 11}),
               ("“기사 끝에 보니까 주택 소유자는 제외라고 딱 적혀 있어요.”",
                {"size": 11})], after=4, lh=1.2)
    para(tfq, [("조유정(31)  ", {"bold": True, "color": BLUE, "size": 11}),
               ("“중위소득 60%면 생각보다 낮거든요. 저는 기준에서 걸려요.”",
                {"size": 11})], after=0, lh=1.2)
    # 우상: 평균 이동 차트(핵심 증거)
    img_fit(slide, ROOT / "eval" / "ablation_shift_viz.png",
            6.7, 1.7, 6.05, 2.55)
    cap = tb(slide, 6.7, 4.28, 6.05, 0.32)
    para(cap, "익명(회색) → 인물 카드(파랑): 응답이 자기 처지로 수렴 — "
         "페르소나가 차이를 만든다", size=10, color=SUB, first=True, after=0,
         align=PP_ALIGN.CENTER)
    # 우하: 행동 벤치마크
    box(slide, 6.7, 4.6, 6.05, 1.8, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf3 = tb(slide, 6.92, 4.74, 5.6, 1.55)
    para(tf3, "③ 강건성 — 행동 벤치마크 7/7", size=12.5, bold=True,
         color=GREEN, first=True, after=5)
    para(tf3, "부조리 가상 정책 5종이 매번 상식적인 집단 축으로 갈라짐 — 반려묘 "
         "유일 찬성자 = ‘펫푸드 창업 꿈’ 페르소나 · 남성 수혜 정책에 남성 23% "
         "반대(“우리 집은 0원”)", size=11, color=INK, after=4, lh=1.25)
    para(tf3, "5종 전부 데모 녹화본(키 없이 재생 가능) · 모델/프롬프트 변경 시 "
         "재실행하는 회귀 도구", size=10, color=SUB, after=0, lh=1.2)
    note(slide, "검증 ②는 저희 발표에서 가장 정직한 대목입니다. 인물 카드 ON/OFF "
         "실험이 사전등록 0/4로 떨어졌고, 점수는 그대로 둔 채 원자료를 부검했더니 "
         "반전이었습니다. 우리 채점 기준은 나이·소득만 보고 10명을 '대상자'라 "
         "표시했는데, 시뮬 속 시민들은 정책 원문의 별도거주·무주택·소득 요건까지 읽고 "
         "전원이 '나는 해당 안 됨'을 정확히 판정한 겁니다 — 실제 청년월세도 신청자 "
         "2/3가 요건 탈락이니, 시뮬이 채점표보다 현실적이었던 거죠. 오른쪽 그림이 "
         "증거입니다. 익명일 땐 '나도 도움되겠지'로 낙관하던 응답이 인물 카드가 "
         "들어가는 순간 자기 처지로 수렴합니다. 검증 ③ 행동 벤치마크는 부조리 가상 "
         "정책 5종을 넣고 사회가 상식 방향으로 갈라지는지 보는 고정 시험지 — 7개 체크 "
         "전부 통과했고, 백미는 전 국민 반려묘 보급의 유일한 찬성자가 페르소나 서사에 "
         "'펫푸드 창업 꿈'이 있는 시민이었다는 것. 입장이 인구통계를 넘어 서사와 "
         "정렬된다는 증거입니다.")


def t09_board(prs):
    """게시판 RAG — 풀버전 18장(번호·노트만 압축)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 10, "05 확장", "게시판 RAG — 문서 근거로 답하는 정책 문의 게시판",
           accent=RED, title_size=24)
    tf = tb(slide, 0.6, 1.8, 5.7, 2.7)
    para(tf, "정책 문의 게시판 + 문서 근거 답변", size=15.5, bold=True,
         color=NAVY, first=True, after=8)
    bullets(tf, [
        ([("질문 → 검색된 근거 ‘안에서만’ 답변", {"bold": True}),
          ("  (환각 억제)", {"color": SUB})], {}),
        ("답변에 근거·품질지표를 함께 표시", {}),
        ("PDF/TXT/MD 업로드로 문서 확장", {}),
        ("API 키가 없으면 추출식 폴백 — 데모는 항상 동작", {}),
    ], size=12.5, after=7, first_done=True)
    box(slide, 0.6, 4.7, 5.7, 1.95, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf2 = tb(slide, 0.85, 4.88, 5.2, 1.6)
    para(tf2, "협업 방식 — 자리만 비워두고 끼워 넣기", size=13.5, bold=True,
         color=GREEN, first=True, after=5)
    para(tf2, "본체에는 끼울 자리와 반환 계약만 미리 설계해 두고, 팀원이 "
         "독립 패키지로 RAG 엔진을 구현해 마지막에 끼워 넣었다. 모듈 경계를 "
         "지킨 덕에 본체는 한 줄도 깨지지 않았다.", size=11.8, color=INK,
         after=0, lh=1.3)
    shot(slide, 6.6, 1.8, 6.13, 4.85, "‘게시판’ 탭 캡처",
         "질문 → 답변 + 📎근거 + 품질지표 표가 보이는 화면 추천")
    note(slide, "확장 파트 첫 번째, 게시판 RAG입니다. 초기에 'MVP에서 제외'로 "
         "좁혔던 RAG가 마지막에 팀의 손으로 돌아왔습니다 — 좁힌 건 버린 게 아니라 "
         "확장 경로였습니다. 시민이나 입안자가 정책 문서에 질문하면 벡터 검색으로 "
         "찾은 근거 '안에서만' 답하고, 근거와 품질지표를 함께 보여줍니다. 협업 방식이 "
         "포인트 — 본체에는 자리와 반환 계약만 설계해 두고 팀원이 독립 패키지로 "
         "구현해 마지막에 끼워 넣었는데, 본체는 한 줄도 깨지지 않았습니다. 키가 "
         "없으면 추출식 폴백으로 동작하니 데모도 안전합니다.")


def t09b_village(prs):
    """미리마을 — 풀버전 19·20장을 한 장으로(마을 + 전파 실측 + 데모 콜아웃)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 11, "05 확장", "미리마을 — 정책 소식은 누구에게, 어떤 경로로 닿는가",
           accent=RED, title_size=23)
    tf = tb(slide, 0.6, 1.75, 6.0, 2.0)
    bullets(tf, [
        ([("Generative Agents", {"bold": True, "color": BLUE}),
          (" (Park et al., 2023) 원형 — LLM이 10명 각자의 하루를 생성, "
           "마주치면 맥락 대화", {})], {}),
        ("손제작 10명 고정 캐스트 — 전파엔 ‘서로 아는 관계’가 필수 "
         "(통계 샘플과 목적이 다른 의도된 선택)", {}),
        ([("녹화 후 재생 — 브라우저는 LLM 0콜", {"bold": True}),
          ("  (즉시·반복·오프라인 = 발표 안전망)", {"color": SUB})], {}),
    ], size=11.8, after=6, lh=1.25)
    box(slide, 0.6, 3.9, 6.0, 1.55, fill=CARD, line=LINE)
    tf2 = tb(slide, 0.85, 4.04, 5.5, 1.3)
    para(tf2, "1일차 실측 (실제 LLM) — 담당 공무원 영희만 알고 시작",
         size=12.5, bold=True, color=NAVY, first=True, after=5)
    para(tf2, [("10명 중 9명에게 전파", {"bold": True, "size": 11.5}),
               (" (공무원→카페 사랑방) · ", {"size": 11.5}),
               ("사각 1명", {"bold": True, "color": RED, "size": 11.5}),
               (" = 동선이 안 겹친 박어르신 · ", {"size": 11.5}),
               ("인지 누수 0", {"bold": True, "color": GREEN, "size": 11.5}),
               (" — ‘아는 사람을 만나야 안다’ 코드 가드", {"size": 11.5})],
         after=0, lh=1.3)
    box(slide, 0.6, 5.6, 6.0, 1.0, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf3 = tb(slide, 0.82, 5.7, 5.6, 0.82, anchor=MSO_ANCHOR.MIDDLE)
    para(tf3, [("창발 —  ", {"bold": True, "color": GREEN, "size": 11.5}),
               ("시키지 않았는데 민수·수아가 광장 산책을 약속하고, 같은 카페 "
                "사장이 손님마다 다른 대화를 한다", {"size": 11})], first=True,
         after=0, lh=1.2)
    # 우: 마을 맵 크게
    img_fit(slide, ROOT / "미리마을" / "assets" / "map.png",
            6.85, 1.75, 5.9, 4.85)
    # 데모 콜아웃
    box(slide, 0.6, 6.78, 12.13, 0.55, fill=PALE_RED, line=RED, lw=1.0)
    tfd = tb(slide, 0.88, 6.82, 11.6, 0.47, anchor=MSO_ANCHOR.MIDDLE)
    para(tfd, "🎬 여기서 라이브 데모 — 미리마을 재생 60~90초 (재생 = LLM 0콜, "
         "네트워크가 죽어도 안전)", size=12, bold=True, color=RED,
         first=True, after=0, align=PP_ALIGN.CENTER)
    note(slide, "마지막 확장, 미리마을입니다. Generative Agents 논문의 원형 — "
         "에이전트가 공간에서 자기 스케줄대로 살아가는 마을입니다. LLM이 10명 각자의 "
         "하루를 생성하면 브라우저는 호출 없이 재생만 하므로 발표 중에도 안전합니다. "
         "축1·2의 통계 샘플엔 '관계'가 없어 전파를 그릴 근거가 없지만, 여긴 서로 아는 "
         "손제작 10명이라 전파가 연출이 아니라 관계와 동선의 결과입니다. 정책 담당 "
         "공무원 영희만 아는 상태로 시작하면 하루 만에 9명에게 퍼지고, 동선이 안 겹친 "
         "박어르신이 사각지대로 남습니다 — 시행착오 장의 '전파 게이트'가 이 무대의 가드입니다. "
         "시키지 않았는데 캐릭터들이 서로를 일정에 엮는 창발도 나타났습니다. 그리고 "
         "발표 첫 장에서 약속한 장면 — 이 친구들이 직접 움직이는 모습을 지금 "
         "보여드리겠습니다. (데모 60~90초. Q&A 대비 — '손으로 안 만든다더니 모순 "
         "아니냐': 통계 샘플은 대표성이 목적, 미리마을은 전파가 목적이라 서로 아는 "
         "관계가 필수 — 목적이 다르면 재료도 다릅니다.)")


def t11_usecases(prs):
    """활용 방안 — 신설(풀버전에 없음). 정책 수명주기 프레임, 실제 탭 기능 근거.

    근거: ui/tab_improve.py(진단 4카드 → 개선안 → 종합 리포트, A/B 재시뮬은
    폐지·수정안 재입력 수동 루프) + 게시판 RAG + 대시보드·인생극장.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 12, "05 활용", "활용 방안 — 정책의 수명주기마다 쓰인다",
           accent=BLUE, title_size=24)
    stages = [
        ("① 입안 — 배포 전 점검", NAVY, PALE_NAVY,
         ["정책 초안을 넣어 사각지대를 미리 본다",
          "· 누가 모르는가 (전파 실패)",
          "· 누가 막히는가 (집행·진입 장벽)",
          "· 악용 위험·허점은 없는가",
          "→ 시민 반응 · 인생극장 · 대시보드"]),
        ("② 다듬기 — 개선 루프", GREEN, PALE_GREEN,
         ["진단 → 개선안 → 보고까지 한 화면",
          "· 병목·사각·우선지원·악용 4카드 진단",
          "· 문구·절차 수정안 + 도움창구 제안",
          "· 수정안을 다시 넣어 재실험(루프)",
          "→ 정책 개선 탭 · 종합 리포트(.md)"]),
        ("③ 시행 후 — 대응 준비", ORANGE, PALE_ORANGE,
         ["민원·문의 대응을 미리 구축한다",
          "· 정책 문서 기반 예상 문의·답변",
          "· 근거 ‘안에서만’ 답변 (환각 억제)",
          "· 안내문·FAQ 작성 재료로 활용",
          "→ 정책 문의 게시판 (RAG)"]),
    ]
    x = 0.6
    for title, c, pale, lines in stages:
        box(slide, x, 1.85, 3.7, 3.6, fill=pale, line=c, lw=1.0)
        tf = tb(slide, x + 0.22, 2.05, 3.3, 3.25)
        para(tf, title, size=14, bold=True, color=c, first=True, after=7)
        para(tf, lines[0], size=11.8, bold=True, color=INK, after=6, lh=1.2)
        for ln in lines[1:-1]:
            para(tf, ln, size=11, color=INK, after=4, lh=1.2)
        para(tf, lines[-1], size=10.5, bold=True, color=c, after=0, lh=1.2)
        x += 4.22
    arrow_text(slide, 4.22, 3.4, 0.5, "→", size=20)
    arrow_text(slide, 8.44, 3.4, 0.5, "→", size=20)
    # 비전 한 줄
    box(slide, 0.6, 5.75, 12.13, 0.95, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tfv = tb(slide, 0.95, 5.87, 11.5, 0.72, anchor=MSO_ANCHOR.MIDDLE)
    para(tfv, [("정책만이 아니다 — ", {"color": INK}),
               ("사내 제도 · 학교 규정 · 공공 캠페인", {"bold": True,
                                                  "color": NAVY}),
               ("처럼 ‘규칙이 사람에게 닿는 모든 곳’이 실험 대상이 될 수 있다",
                {"color": INK})], size=14, first=True, after=0,
         align=PP_ALIGN.CENTER)
    note(slide, "그래서 누가, 어떻게 쓰느냐 — 정책의 수명주기를 따라갑니다. 입안 "
         "단계에선 초안을 넣어 누가 모르고, 누가 막히고, 어떤 악용 위험이 있는지 "
         "배포 전에 봅니다. 다듬기 단계에선 정책 개선 탭이 병목·사각·우선지원·악용 "
         "네 카드로 진단하고, 문구·절차 수정안과 도움창구 운영 제안을 만들어 "
         "종합 리포트로 내보냅니다 — 수정안을 다시 정책 입력에 넣으면 개선 루프가 "
         "돌아갑니다. 시행 후엔 게시판 RAG로 정책 문서 기반 예상 문의·답변을 미리 "
         "구축해 민원 대응을 준비합니다. 그리고 이 구조는 정책 전용이 아닙니다 — "
         "사내 제도, 학교 규정처럼 규칙이 사람에게 닿는 모든 곳에 같은 방식으로 쓸 "
         "수 있습니다.")


def t10_limits(prs):
    """한계 — 풀버전 22장의 5가지 중 핵심 3가지."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 13, "06 마무리", "한계 — 먼저 말해두는 세 가지", accent=SUB)
    items = [
        ("미래를 맞히는 도구가 아니다",
         "일관되게 ‘영향 시나리오’로 포지셔닝. 신뢰성은 예측 정확도가 아니라 "
         "검증 3축(재료·작동·강건성)과 사전등록으로 방어한다."),
        ("표본이 작다",
         "시민 24명 — 비율 최소 단위 4.2%p. 통계적 일반화가 아니라 ‘구조가 "
         "나타나는지’의 시연이다."),
        ("이념·가치 축이 없다",
         "과거 정책 갭 측정(시뮬 vs 당시 여론조사)에서 찬성률 MAE 27.4%p — 주범은 "
         "반대 진영 소실. 처방은 프롬프트가 아니라 페르소나 데이터 강화."),
    ]
    y = 1.95
    for title, body in items:
        box(slide, 0.6, y, 12.13, 1.25, fill=CARD, line=LINE)
        tf = tb(slide, 0.88, y + 0.14, 11.6, 1.0)
        para(tf, title, size=14.5, bold=True, color=NAVY, first=True, after=4)
        para(tf, body, size=12.5, color=INK, after=0, lh=1.25)
        y += 1.45
    f = tb(slide, 0.6, 6.45, 12.13, 0.4)
    para(f, "그 외(소득 판정 근사 · ‘실현 기대’ 측정축 부재)는 자료집 22장에 정리",
         size=10.5, color=SUB, first=True, after=0)
    note(slide, "마무리 전에 한계를 저희 입으로 먼저 말씀드립니다. 미리랩은 미래를 "
         "맞히는 도구가 아니고 그렇게 주장하지도 않습니다 — 영향 시나리오입니다. "
         "표본 24명은 통계적 일반화가 아니라 구조가 나타나는지의 시연입니다. 그리고 "
         "가장 큰 공백을 정직하게 — 과거 정책 4건을 당시 여론조사와 맞대본 갭 측정에서 "
         "찬성률 평균 오차 27.4%p였고, 주범은 이념·가치 축 부재로 인한 '반대 진영 "
         "소실'입니다. 프롬프트 수정이 아니라 페르소나 데이터 강화로 풀 문제로 진단해 "
         "뒀습니다. 이 정직한 선 긋기가 저희 설계의 일부입니다.")


# ---------------------------------------------------------------- 메인
def main():
    sprites = prep_sprites()
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    s01_title(prs, sprites)      # 1  타이틀 (풀버전과 공유)
    t02_problem_concept(prs)     # 2  문제 + 컨셉
    t03_system(prs)              # 3  시스템 한 장
    t04_axis1(prs)               # 4  축1 (캡처)
    t05_axis2(prs)               # 5  축2 (캡처)
    t06_axis3(prs)               # 6  축3 (캡처)
    t07_lessons(prs)             # 7  시행착오 대표 2가지
    t08_validation(prs)          # 8  검증 개요 + ① 재료
    t08b_audit(prs)              # 9  검증 ②③ — 부검 반전 + 행동 벤치
    t09_board(prs)               # 10 게시판 RAG (캡처)
    t09b_village(prs)            # 11 미리마을 + 데모 콜아웃
    t11_usecases(prs)            # 12 활용 방안 (신설)
    t10_limits(prs)              # 13 한계 3가지
    s19_closing(prs, sprites)    # 14 클로징 (풀버전과 공유)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print("[OK] saved:", OUT)
    print("[OK] slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
