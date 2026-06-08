# -*- coding: utf-8 -*-
"""미리랩 발표자료(.pptx) 생성 스크립트.

사용법:  python _make_pptx.py
출력:    notebooks/미리랩_발표.pptx  (21슬라이드, 발표자 노트 포함)

- 내용 출처: notebooks/발표_구현정리.md + 검증 3축(eval/persona_eval·ablation·robustness)
  + 행동 벤치마크 v1(eval/behavior_bench_*) — 2026-06-07 검증 세션 반영
- 임베드 차트: persona_eval_viz / ablation_shift_viz / behavior_bench_viz
  (뒤 2장은 _make_bench_viz.py 가 eval JSON 에서 생성 — 검증 재실행 시 같이 재실행)
- 스크린샷 자리는 회색 점선 박스(📸)로 비워둠 — 캡처 후 교체
- 다시 만들려면 이 스크립트만 재실행 (기존 pptx 덮어씀)
"""
from __future__ import annotations

import json
import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

try:  # 점선 테두리 (없으면 실선으로 폴백)
    from pptx.enum.line import MSO_LINE

    DASH = MSO_LINE.DASH
except Exception:  # pragma: no cover
    DASH = None

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "notebooks" / "미리랩_발표.pptx"
ASSETS = ROOT / "notebooks" / "_pptx_assets"

# ---------------------------------------------------------------- 팔레트
NAVY = RGBColor(0x1F, 0x3A, 0x5F)      # 본문 제목/메인
BLUE = RGBColor(0x2D, 0x6C, 0xDF)      # 포인트
INK = RGBColor(0x1F, 0x29, 0x37)       # 본문 텍스트
SUB = RGBColor(0x6B, 0x72, 0x80)       # 보조 텍스트
GREEN = RGBColor(0x2E, 0x9E, 0x6B)     # 수혜
ORANGE = RGBColor(0xE8, 0x90, 0x2A)    # 경계
RED = RGBColor(0xD9, 0x53, 0x4F)       # 사각지대
CARD = RGBColor(0xF4, 0xF6, 0xFA)      # 카드 배경
LINE = RGBColor(0xDD, 0xE3, 0xEC)      # 카드 테두리
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE_NAVY = RGBColor(0xEA, 0xF0, 0xF8)
PALE_GREEN = RGBColor(0xEA, 0xF6, 0xF0)
PALE_ORANGE = RGBColor(0xFD, 0xF3, 0xE4)
PALE_RED = RGBColor(0xFC, 0xEC, 0xEB)
GRAY_BG = RGBColor(0xF1, 0xF3, 0xF7)
GRAY_BORDER = RGBColor(0x9A, 0xA5, 0xB1)

FONT = "맑은 고딕"
SW, SH = 13.333, 7.5  # 슬라이드 크기(인치)


# ---------------------------------------------------------------- 헬퍼
def _font(run, *, size=14, bold=False, color=INK, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def tb(slide, x, y, w, h, *, anchor=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def para(tf, content, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
         after=4, before=0, lh=1.12, first=False, italic=False):
    """content: str 또는 [(str, {스타일 오버라이드}), ...]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    p.line_spacing = lh
    if isinstance(content, str):
        content = [(content, {})]
    for text, ov in content:
        r = p.add_run()
        r.text = text
        _font(r, size=ov.get("size", size), bold=ov.get("bold", bold),
              color=ov.get("color", color), italic=ov.get("italic", italic))
    return p


def box(slide, x, y, w, h, *, fill=CARD, line=LINE, lw=0.75, dash=False,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
        if dash and DASH is not None:
            s.line.dash_style = DASH
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    s.text_frame.margin_left = s.text_frame.margin_right = Inches(0.12)
    s.text_frame.margin_top = s.text_frame.margin_bottom = Inches(0.06)
    return s


def chip(slide, x, y, text, *, fill=NAVY, color=WHITE, size=11, w=1.6, h=0.3):
    c = box(slide, x, y, w, h, fill=fill, line=None, radius=0.5)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, text, size=size, bold=True, color=color, align=PP_ALIGN.CENTER,
         after=0, first=True)
    return c


def header(slide, no, sec, title, *, accent=NAVY, title_size=27):
    chip(slide, 0.6, 0.42, sec, fill=accent, w=1.75, h=0.32)
    t = tb(slide, 0.6, 0.82, SW - 1.2, 0.75)
    para(t, title, size=title_size, bold=True, color=NAVY, first=True, after=0)
    # 푸터
    f = tb(slide, 0.6, SH - 0.36, SW - 1.2, 0.3)
    para(f, [("미리랩", {"bold": True, "color": SUB}),
             ("  ·  AI 시민 사회 정책 실험실", {"color": SUB})],
         size=9, first=True, after=0)
    n = tb(slide, SW - 1.2, SH - 0.36, 0.6, 0.3)
    para(n, str(no), size=9, color=SUB, align=PP_ALIGN.RIGHT, first=True, after=0)


def shot(slide, x, y, w, h, label, sub="앱에서 캡처한 뒤, 이 박스를 지우고 이미지를 넣어주세요"):
    s = box(slide, x, y, w, h, fill=GRAY_BG, line=GRAY_BORDER, lw=1.2,
            dash=True, radius=0.04)
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "📸 " + label, size=14, bold=True, color=SUB,
         align=PP_ALIGN.CENTER, first=True, after=2)
    para(tf, sub, size=10, color=GRAY_BORDER, align=PP_ALIGN.CENTER, after=0)
    return s


def img_fit(slide, path, x, y, w, h):
    """비율 유지로 (x,y,w,h) 박스 안에 이미지 배치(가운데 정렬)."""
    iw, ih = Image.open(path).size
    box_ratio, img_ratio = w / h, iw / ih
    if img_ratio >= box_ratio:
        dw, dh = w, w / img_ratio
    else:
        dh, dw = h, h * img_ratio
    px = x + (w - dw) / 2
    py = y + (h - dh) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py),
                                    Inches(dw), Inches(dh))


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def role_card(slide, x, y, w, h, color, pale, title, lines, *, title_size=14,
              body_size=11.5):
    box(slide, x, y, w, h, fill=pale, line=color, lw=1.0)
    box(slide, x, y, 0.07, h, fill=color, line=None, radius=0.5)
    tf = tb(slide, x + 0.2, y + 0.12, w - 0.35, h - 0.24)
    para(tf, title, size=title_size, bold=True, color=color, first=True, after=4)
    for ln in lines:
        para(tf, ln, size=body_size, color=INK, after=2, lh=1.15)


def flow_box(slide, x, y, w, h, title, body, *, fill=PALE_NAVY, line_c=NAVY,
             tcolor=NAVY, tsize=12.5, bsize=10.5):
    box(slide, x, y, w, h, fill=fill, line=line_c, lw=1.0)
    tf = tb(slide, x + 0.12, y + 0.08, w - 0.24, h - 0.16)
    para(tf, title, size=tsize, bold=True, color=tcolor, first=True, after=2)
    if body:
        for ln in body if isinstance(body, list) else [body]:
            para(tf, ln, size=bsize, color=INK, after=1, lh=1.12)


def arrow_text(slide, x, y, w, ch="→", *, size=18, color=SUB):
    t = tb(slide, x, y, w, 0.4)
    para(t, ch, size=size, bold=True, color=color, align=PP_ALIGN.CENTER,
         first=True, after=0)


def bullets(tf, items, *, size=13, after=5, lh=1.18, first_done=False):
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            text, ov = it
        else:
            text, ov = it, {}
        content = [("•  ", {"color": ov.get("bullet_color", BLUE), "bold": True})]
        if isinstance(text, list):  # 이미 (run, style) 리스트인 경우
            content.extend(text)
        else:
            content.append((text, ov))
        para(tf, content, size=ov.get("size", size), after=ov.get("after", after),
             lh=lh, first=(i == 0 and not first_done))


# ---------------------------------------------------------------- 스프라이트 준비
def prep_sprites():
    ASSETS.mkdir(exist_ok=True)
    defs = json.loads((ROOT / "미리마을" / "assets" / "sprite_defs.json")
                      .read_text(encoding="utf-8"))
    out = {}
    for key, d in defs.items():
        sheet = ROOT / "미리마을" / d["sheet"].replace("/", os.sep)
        if not sheet.exists():
            continue
        cell = d.get("cell", 48)
        img = Image.open(sheet).convert("RGBA")
        # 정면(down) 행의 가운데 프레임
        row = d.get("rowMap", {}).get("down", 0)
        frame = img.crop((cell, row * cell, cell * 2, (row + 1) * cell))
        if frame.getbbox() is None:  # 빈 프레임이면 첫 프레임
            frame = img.crop((0, row * cell, cell, (row + 1) * cell))
        frame = frame.resize((cell * 4, cell * 4), Image.NEAREST)
        p = ASSETS / f"sprite_{key}.png"
        frame.save(p)
        out[key] = p
    return out


# ================================================================ 슬라이드들
def s01_title(prs, sprites):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 상단 컬러 바
    box(slide, 0, 0, SW, 0.18, fill=NAVY, line=None,
        shape=MSO_SHAPE.RECTANGLE)
    t = tb(slide, 1.0, 1.55, SW - 2.0, 3.6)
    para(t, "미리랩", size=64, bold=True, color=NAVY, first=True, after=2)
    para(t, "MiriLab — AI 시민 사회 정책 실험실", size=20, color=BLUE,
         bold=True, after=18)
    para(t, "정책을 배포하기 전에,", size=22, color=INK, after=2)
    para(t, [("실제 한국 인구통계 기반 ", {}),
             ("AI 가상 시민", {"bold": True, "color": BLUE}),
             ("에게 먼저 물어본다.", {})], size=22, color=INK, after=16)
    para(t, [("수혜 ", {"bold": True, "color": GREEN}),
             ("·  ", {"color": SUB}),
             ("경계 ", {"bold": True, "color": ORANGE}),
             ("·  ", {"color": SUB}),
             ("사각지대 ", {"bold": True, "color": RED}),
             ("— 같은 정책, 다른 인생", {"color": SUB})], size=15, after=0)
    # 하단: 스프라이트 행렬
    keys = ["minsu", "sua", "owner", "staff", "grandma", "oldman", "miyoung",
            "junho", "jimin", "daeun"]
    x = 1.0
    for k in keys:
        p = sprites.get(k)
        if p:
            slide.shapes.add_picture(str(p), Inches(x), Inches(5.6),
                                     Inches(0.62), Inches(0.62))
            x += 0.78
    info = tb(slide, 1.0, 6.55, SW - 2.0, 0.6)
    para(info, "DLthon 2  ·  팀명/발표자 (채워주세요)  ·  2026. 06.",
         size=13, color=SUB, first=True, after=0)
    note(slide, "안녕하세요, 저희는 '미리랩'을 만들었습니다. "
         "미리랩은 정책을 배포하기 '전에' AI 가상 시민 사회에서 '미리' 반응을 실험하는 "
         "랩(Lab)입니다. 아래 캐릭터들은 실제로 저희 시뮬레이션 속에서 살아가는 "
         "가상 시민들인데, 발표 마지막에 이 친구들이 마을에서 직접 움직이는 모습을 "
         "보여드리겠습니다. (팀명·발표자 이름은 슬라이드에 채워주세요)")


def s02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 2, "발표 흐름", "발표 순서 (약 12~15분)")
    rows = [
        ("01", "문제와 컨셉", "정책은 왜 시민에게 닿지 못하는가 → 미리랩의 답", "약 2분", NAVY),
        ("02", "시스템", "페르소나 파이프라인 + 두 축(시민 반응 / 정책 인생극장)", "약 3분", BLUE),
        ("03", "시행착오", "스코프 조정과 시행착오 3가지", "약 3분", ORANGE),
        ("04", "신뢰성 검증", "검증 3축(재료·작동·한계) + 사전등록 + 행동 벤치마크", "약 3분", GREEN),
        ("05", "확장 구현", "게시판 RAG + 미리마을 + 60초 라이브 재생 데모", "약 3분", RED),
        ("06", "한계와 마무리", "팀 규율 · 한계 · 핵심 메시지 · Q&A", "약 2분", SUB),
    ]
    y = 1.75
    for no, sec, desc, time, color in rows:
        box(slide, 0.6, y, 12.13, 0.74, fill=CARD, line=LINE)
        c = box(slide, 0.78, y + 0.13, 0.48, 0.48, fill=color, line=None,
                shape=MSO_SHAPE.OVAL)
        tfc = c.text_frame
        tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tfc, no, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             first=True, after=0)
        t = tb(slide, 1.5, y + 0.09, 9.4, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        para(t, [(sec + "   ", {"bold": True, "size": 15, "color": NAVY}),
                 (desc, {"color": SUB, "size": 12.5})], first=True, after=0)
        tt = tb(slide, 11.0, y + 0.09, 1.55, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        para(tt, time, size=12, bold=True, color=color, align=PP_ALIGN.RIGHT,
             first=True, after=0)
        y += 0.86
    note(slide, "발표는 여섯 부분입니다. 문제의식에서 출발해, 시스템 두 축을 보여드리고, "
         "이 모양에 도달하기까지의 시행착오를 솔직하게 공유합니다. 그리고 '이걸 믿을 수 "
         "있나'에 대한 정량 검증, 마지막으로 확장 구현(게시판 RAG와 미리마을)과 한계를 "
         "말씀드립니다. — 섹션 단위로 발표자를 나누기 좋게 구성했습니다. 10분 버전이 "
         "필요하면 02·04·05를 한 장씩 줄이는 걸 추천합니다.")


def s03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 3, "01 문제", "정책은 왜 시민에게 닿지 못하는가")
    cards = [
        ("정보 격차", "내가 대상자인지\n조차 모른다"),
        ("이해도 부족", "문구가 어렵고\n신청 조건이 불명확하다"),
        ("디지털 장벽", "시니어에겐 온라인 신청\n자체가 벽이다"),
        ("도달 실패", "발표돼도 닿지 않는다\n— “그림의 떡”"),
    ]
    x = 0.6
    for title, body in cards:
        box(slide, x, 1.85, 2.93, 1.75, fill=PALE_RED, line=RED, lw=1.0)
        tf = tb(slide, x + 0.18, 2.05, 2.6, 1.4)
        para(tf, title, size=16, bold=True, color=RED, first=True, after=6)
        for ln in body.split("\n"):
            para(tf, ln, size=12.5, color=INK, after=1, lh=1.2)
        x += 3.07
    # 하단 핵심 대비
    box(slide, 0.6, 4.0, 12.13, 2.6, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tf = tb(slide, 0.95, 4.25, 11.5, 2.2)
    para(tf, "기존 서비스의 한계", size=15, bold=True, color=NAVY, first=True,
         after=8)
    para(tf, [("정부24 챗봇 · 복지로  →  정책을 ", {}),
              ("‘찾는’ 사람", {"bold": True, "color": BLUE}),
              ("을 돕는다 (검색·매칭)", {})], size=15, after=6)
    para(tf, [("그러나 정책을 ", {}),
              ("‘만드는’ 사람(입안자)", {"bold": True, "color": RED}),
              ("이 “이 정책, 배포하면 누가 못 받을까?”를", {})], size=15,
         after=2)
    para(tf, "배포 전에 미리 볼 수 있는 도구는 없다.", size=15, bold=True, after=0)
    note(slide, "정부와 지자체는 매년 수천 건의 정책을 발표하지만, 정작 필요한 사람에게 "
         "닿지 못하는 구조적 누수가 반복됩니다. 모르거나, 어렵거나, 디지털이 벽이거나, "
         "발표돼도 닿지 않거나. 기존 서비스는 전부 정책을 '찾는' 시민을 돕는 도구입니다. "
         "정책을 '만드는' 입안자가 배포 전에 '누가 못 받을지'를 미리 보는 도구는 없습니다. "
         "저희는 바로 그 빈자리를 노렸습니다.")


def s04_concept(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 4, "01 컨셉", "미리랩 — 정책을 시민 입장에서 미리 실험한다")
    box(slide, 0.6, 1.75, 12.13, 1.05, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tf = tb(slide, 0.95, 1.95, 11.5, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("“누가 이해 못 하고 · 누가 신청 못 하고 · 어디서 막히는지”",
               {"bold": True, "color": NAVY}),
              ("를 배포 전에 보여준다", {})], size=17, first=True, after=0,
         align=PP_ALIGN.CENTER)
    # 핵심 장면: 같은 정책, 다른 인생
    t = tb(slide, 0.6, 3.05, 12.13, 0.45)
    para(t, [("핵심 장면 — 같은 ‘청년 월세 지원’ 정책, 6개월 후",
              {"bold": True, "size": 15, "color": INK})], first=True, after=0)
    role_card(slide, 0.6, 3.6, 3.93, 2.3, GREEN, PALE_GREEN, "수혜  ·  A (26세 청년)",
              ["복지로에서 5분 만에 신청,", "월세 부담이 줄어", "첫 저축을 시작했다."],
              body_size=13)
    role_card(slide, 4.7, 3.6, 3.93, 2.3, ORANGE, PALE_ORANGE, "경계  ·  B (46세 중장년)",
              ["본인은 대상이 아니지만", "자녀를 위해 알아보고", "대신 신청해 준다."],
              body_size=13)
    role_card(slide, 8.8, 3.6, 3.93, 2.3, RED, PALE_RED, "사각지대  ·  C (74세 고령)",
              ["6개월이 지나도", "정책의 존재조차", "알지 못했다."], body_size=13)
    t2 = tb(slide, 0.6, 6.15, 12.13, 0.6)
    para(t2, "“같은 정책, 다른 인생.” — 미리랩은 이 갈림을 배포 전에 보여준다",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, first=True, after=0)
    note(slide, "미리랩의 한 줄 포지셔닝입니다. 같은 청년 월세 정책이라도 26세 청년은 "
         "5분 만에 신청해 혜택을 받고, 46세 중장년은 자녀를 위해 대신 신청해 주고, "
         "74세 어르신은 6개월이 지나도 존재조차 모릅니다. '같은 정책, 다른 인생' — "
         "이 갈림을 배포 전에 미리 보여주는 것이 미리랩입니다. 이 세 가지 색(수혜·경계·"
         "사각)이 오늘 발표 전체를 관통하는 색입니다.")


def s05_system(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 5, "02 시스템", "시스템 — 단일 LangGraph에서 3축 단방향으로",
           title_size=23)
    # v1 (폐기) — 단일 LangGraph 가 전부였던 시절
    flow_box(slide, 0.6, 1.5, 12.13, 0.8,
             "v1 · 단일 LangGraph:  react → interact → aggregate",
             "한 호출이 인지·점수·결과까지 다 판정  →  출력 과밀 · 축 간 모순",
             fill=GRAY_BG, line_c=GRAY_BORDER, tcolor=SUB, tsize=14, bsize=11.5)
    # 전환
    arrow_text(slide, 5.55, 2.32, 0.5, "↓")
    t = tb(slide, 6.1, 2.38, 6.2, 0.35)
    para(t, "재설계 — 책임을 쪼개 단방향으로", size=12, bold=True, color=BLUE,
         first=True, after=0)
    # v2 컨테이너 (3축 단방향)
    box(slide, 0.6, 2.85, 12.13, 3.95, fill=None, line=GREEN, lw=1.0, dash=True,
        radius=0.03)
    cap = tb(slide, 0.85, 2.92, 6.0, 0.32)
    para(cap, "v2 · 3축 단방향 파이프라인", size=13, bold=True, color=GREEN,
         first=True, after=0)
    # 페르소나 (바닥 재료)
    flow_box(slide, 1.0, 3.35, 11.33, 0.6,
             "페르소나 — Nemotron-Personas-Korea (실제 한국 통계 기반)",
             None, fill=PALE_NAVY, line_c=NAVY, tsize=12.5)
    arrow_text(slide, 6.42, 3.97, 0.5, "↓", size=15)
    # 3축 (단방향)
    flow_box(slide, 1.0, 4.4, 3.5, 1.0, "축1 · 정보",
             "react — t0 반응", fill=PALE_GREEN, line_c=GREEN, tcolor=GREEN,
             tsize=14, bsize=12)
    flow_box(slide, 4.92, 4.4, 3.5, 1.0, "축2 · 결과",
             "인생극장 · 대조 3명", fill=PALE_ORANGE, line_c=ORANGE,
             tcolor=ORANGE, tsize=14, bsize=12)
    flow_box(slide, 8.83, 4.4, 3.5, 1.0, "축3 · 요약",
             "집계·서술 (axis3)", fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY,
             tsize=14, bsize=12)
    arrow_text(slide, 4.42, 4.75, 0.5, "→")
    arrow_text(slide, 8.33, 4.75, 0.5, "→")
    # 정정 주석
    n = tb(slide, 1.0, 5.52, 11.0, 0.3)
    para(n, "LangGraph는 이제 축1만 담당 · 게이지 포함 모든 집계는 축3로 단일화",
         size=11.5, color=SUB, first=True, after=0)
    arrow_text(slide, 6.42, 5.78, 0.5, "↓", size=15)
    # 탭
    flow_box(slide, 1.0, 6.15, 11.33, 0.55,
             "Streamlit 7탭 · 데모 = Gemini 시뮬 녹화 재생 (네트워크 0콜)",
             None, fill=CARD, line_c=NAVY, tsize=12.5)
    note(slide, "이 시스템은 처음부터 이 모양이 아니었습니다. 처음엔 단일 LangGraph "
         "파이프라인 하나가 전부였습니다 — react가 시민 반응을, interact가 전파를, "
         "aggregate가 집계까지 한 흐름에서 다 했습니다. 그런데 한 번의 호출이 "
         "인지·감정·점수·결과 예측까지 너무 많은 걸 판정하다 보니 두 가지가 터졌습니다. "
         "하나는 출력 스키마가 자꾸 깨졌고, 또 하나는 같은 사람의 신청·수령을 여러 곳이 "
         "따로 판정해 서로 모순됐습니다. 그래서 책임을 쪼갰습니다. 축1(정보)은 시민의 "
         "첫 반응만 만들고, 축2(결과)가 그걸 시딩 삼아 시간 속에서 인생을 굴리고, "
         "축3(요약)가 두 축을 읽어 집계하고 서술합니다 — 단방향이라 뒤로 흐르지 않고, "
         "판정 지점이 하나라 모순이 안 생깁니다. LangGraph는 이제 전체가 아니라 축1만 "
         "담당하고, 게이지를 포함한 모든 화면 숫자는 축3 한 곳에서 나옵니다. 바닥엔 실제 "
         "한국 통계 기반 페르소나가 있고, 위는 Streamlit 7탭으로 통합됩니다. 데모는 "
         "Gemini 시뮬을 미리 녹화해 재생하는 방식이라 네트워크가 죽어도 진짜 결과로 "
         "시연됩니다. 시민 모델은 사이드바에서 선택 가능, 현재 기본은 gemini-3-flash.")


def s05b_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 6, "02 시스템", "페르소나 — 어떻게 만들어진 데이터인가")
    tf = tb(slide, 0.6, 1.55, 12.13, 0.5)
    para(tf, "NVIDIA가 실제 통계에서 뽑아 LLM으로 2단계로 빚어 공개한 데이터셋 — 우리는 이걸 가져다 쓴다",
         size=14, bold=True, color=NAVY, first=True, after=0)
    # 3단 생성 흐름 (NVIDIA 가 데이터셋을 만든 과정)
    nlab = tb(slide, 0.6, 2.28, 12.13, 0.28)
    para(nlab, "NVIDIA가 데이터셋을 만든 과정 (우리 작업이 아니라 데이터 출처)",
         size=11, color=SUB, first=True, after=0)
    flow_box(slide, 0.6, 2.62, 3.6, 1.7, "① 실제 분포에서 뽑기",
             ["인구통계: 나이·성별·지역·", "직업·학력·가구형태",
              "+ 성격 5요인을 입힘"],
             fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY, tsize=13.5, bsize=11.5)
    arrow_text(slide, 4.25, 3.25, 0.5, "→")
    flow_box(slide, 4.75, 2.62, 3.6, 1.7, "② LLM 1차 — 살 붙이기",
             ["문화적 배경 · 전문성", "인생 목표 · 취미", "(4가지 특성 생성)"],
             fill=PALE_GREEN, line_c=GREEN, tcolor=GREEN, tsize=13.5, bsize=11.5)
    arrow_text(slide, 8.4, 3.25, 0.5, "→")
    flow_box(slide, 8.9, 2.62, 3.83, 1.7, "③ LLM 2차 — 종합",
             ["직업·가족·음식·여행 등", "영역별 인물 이야기 완성", "(7가지 서사)"],
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=13.5, bsize=11.5)
    # 결과 + 핵심 메시지
    box(slide, 0.6, 4.65, 12.13, 1.7, fill=CARD, line=LINE)
    tf2 = tb(slide, 0.85, 4.8, 11.6, 1.45)
    para(tf2, [("결과 — ", {"bold": True, "color": NAVY}),
               ("100만 명 × 26개 항목", {"bold": True, "color": BLUE}),
               ("  (인구통계 12개 + LLM이 쓴 인물 서사 14개)", {"color": INK})],
         size=14, first=True, after=6)
    para(tf2, "위 ①②③의 특성과 인물 서사가 데이터셋 컬럼에 그대로 들어 있다 — "
         "도식이 아니라 실제 데이터.",
         size=12, color=SUB, after=8, lh=1.25)
    para(tf2, [("→ 통계적으로 실재하고 성격으로 일관된 시민. ",
                {"bold": True, "color": GREEN}),
               ("사각지대 인물이 처음부터 표본 안에 들어 있다.",
                {"bold": True, "color": INK})],
         size=13, after=0, lh=1.25)
    note(slide, "이 페르소나가 어떻게 만들어졌는지가 신뢰의 출발점입니다. NVIDIA의 "
         "NeMo Data Designer가 두 단계로 빚습니다. 먼저 실제 한국 인구통계 분포에서 "
         "나이·성별·지역·직업·학력·가구형태를 뽑고, 성격 5요인 모형으로 성격의 결을 "
         "입힙니다 — 그래서 '76세·초졸·무직·하역종사원'처럼 통계적으로 함께 나타날 법한 "
         "조합이 자연스럽게 나옵니다. 그다음 LLM이 1차로 그 사람의 문화적 배경·전문성·"
         "인생 목표·취미를 채우고, 2차로 이 모두를 종합해 직업·가족·음식·여행 같은 영역별 "
         "인물 이야기를 완성합니다. 결과는 100만 명 × 26개 항목입니다. 핵심은 — 우리가 "
         "보고 싶은 반응을 손으로 심은 게 아니라 실제 분포에서 뽑았기 때문에, 정책 "
         "사각지대에 놓일 인물이 처음부터 표본 안에 들어 있다는 점입니다. 이게 다음 장의 "
         "'우리가 이 데이터를 어떻게 쓰는가'의 전제가 됩니다.")


def s06_persona(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 7, "02 시스템", "페르소나 — 손으로 만들지 않고 통계에서 뽑는다")
    tf = tb(slide, 0.6, 1.75, 5.9, 4.6)
    para(tf, "원칙: 그럴듯한 가짜 대신 실제 통계", size=16, bold=True,
         color=NAVY, first=True, after=8)
    bullets(tf, [
        ([("nvidia/Nemotron-Personas-Korea", {"bold": True, "color": BLUE}),
          ("  — 통계청·대법원·건보공단 기반 합성 데이터셋", {})], {}),
        ("약 100만 행 × 26컬럼, CC BY 4.0 (상업 활용 가능)", {}),
        ("첫 샤드(220MB)만 + seed 고정 샘플링 + 로컬 캐시", {}),
        ([("→ 재현성 ", {"bold": True}),
          ("(발표마다 같은 시민, 오프라인에서도 동작)", {"color": SUB})], {}),
    ], size=13.5, after=7, first_done=True)
    para(tf, "", size=6, after=2)
    para(tf, "왜 중요한가", size=15, bold=True, color=NAVY, after=6)
    para(tf, "개요 초안은 “페르소나 20~30명 직접 설계”였다. 손으로 만들면 "
         "우리가 보고 싶은 반응을 심게 된다. 실제 분포에서 샘플링하면 "
         "사각지대 인물이 처음부터 표본 안에 들어 있다.",
         size=13, color=INK, after=0, lh=1.3)
    # 우측: 결정론 신호 표
    bx, by, bw = 6.85, 1.75, 5.88
    box(slide, bx, by, bw, 4.6, fill=CARD, line=LINE)
    tf2 = tb(slide, bx + 0.25, by + 0.2, bw - 0.5, 4.2)
    para(tf2, "기존 인구통계에서 규칙으로 파생한 보조 신호 (LLM 없이 · 같은 사람이면 늘 같은 값)",
         size=13.5, bold=True, color=NAVY, first=True, after=9)
    para(tf2, [("digital_literacy", {"bold": True, "color": BLUE, "size": 13})], after=2)
    para(tf2, "나이·학력·직업으로 계산  →  인물 카드 입력 · 접근도 · 검증에 사용 "
         "(디지털 장벽 = 사각지대의 핵심)", size=11.5, color=INK, after=9, lh=1.2)
    para(tf2, [("income_level", {"bold": True, "color": BLUE, "size": 13})], after=2)
    para(tf2, "직업·학력 키워드로 분류  →  정책 대상 판정(누가 수혜 대상인가)에 사용",
         size=11.5, color=INK, after=9, lh=1.2)
    para(tf2, [("그 외 보조 2종 — ", {"bold": True, "color": SUB, "size": 11.5}),
               ("government_trust(접근도 가중 0.25) · social_network(전파용). "
                "발표 라인엔 거의 미사용.", {"size": 11.5, "color": SUB})],
         after=9, lh=1.2)
    para(tf2, [("검증 ① 통과 — ", {"bold": True, "color": GREEN, "size": 11.5}),
               ("표본 24명, 나이·성별·지역·학력 등 7개 변수 전부 무작위 추출 "
                "기대범위 안 (부트스트랩 1만회)", {"size": 11.5, "color": INK})],
         after=0, lh=1.2)
    note(slide, "페르소나는 손으로 만들지 않았습니다. NVIDIA가 통계청·대법원·건보공단 등 "
         "실제 통계로 합성한 100만 행 데이터셋에서 seed 고정으로 샘플링합니다. 손으로 "
         "만들면 보고 싶은 반응을 심게 되니까요. 실제 분포에서 뽑으면 사각지대 인물이 "
         "'이미 시드 안에' 들어 있습니다. 데이터셋 자체의 현실 정합은 제작사(NVIDIA)가 "
         "카드에 검증·문서화했고(고령층 성비 1.52배, 최다 성씨 김 21.5% 등), 저희는 "
         "'거기서 쏠림 없이 뽑았는가'를 부트스트랩 1만회로 검증해 7개 변수 전부 "
         "통과했습니다(검증 ①, 뒤 04장). 데이터에 없는 신호는 순수 파이썬 규칙으로 "
         "유도하고, 정부 신뢰처럼 정보가 없는 값은 자리표시자임을 정직하게 표기합니다. "
         "재현성이 발표 안정성이기도 합니다.")


def s07_axis_a(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 8, "02 시스템", "축1 · 정보 — 시민 반응", accent=GREEN)
    fx, fw = 0.6, 5.8
    para(tb(slide, fx, 1.75, fw, 0.6),
         "시민 각자가 정책을 보고 즉시 반응한다 (t0)", size=15, bold=True,
         color=NAVY, first=True, after=0)
    flow_box(slide, fx, 2.5, fw, 1.55, "react — 1차 반응",
             ["· 반응문 + 입장 (찬성 · 반대 · 혼합)",
              "· 익명 설문 응답",
              "· LLM은 선택지만 고르고, 0~100 점수 변환은 코드"],
             fill=PALE_GREEN, line_c=GREEN, tcolor=GREEN, tsize=14, bsize=12)
    box(slide, fx, 4.32, fw, 1.05, fill=GRAY_BG, line=GRAY_BORDER, dash=True)
    tfi = tb(slide, fx + 0.22, 4.43, fw - 0.44, 0.85)
    para(tfi, "interact — 다른 시민을 보고 입장 변화", size=12.5, bold=True,
         color=SUB, first=True, after=2)
    para(tfi, "SNS 채팅 탭 데모 전용 · 집계엔 미반영 (메인 라인에서 덜어냄)",
         size=11.5, color=SUB, after=0)
    para(tb(slide, fx, 5.65, fw, 0.5),
         "→ 집단의 첫 온도. 게이지·리포트 집계는 축3에서.", size=12.5,
         bold=True, color=GREEN, first=True, after=0)
    para(tb(slide, fx, 6.42, fw, 0.45),
         "캐시 — 다중 시민 호출이 공통 프롬프트(prefix)를 공유 (cached_tokens 실측)",
         size=10.5, color=SUB, first=True, after=0)
    shot(slide, 6.7, 1.75, 6.03, 5.05, "‘시민 반응’ 탭 캡처",
         "데모에서 실제 반응이 채워지는 화면")
    note(slide, "축1은 시민 각자의 첫 반응(t0)을 만드는 단계입니다. 정책을 보고 "
         "반응문과 입장(찬성·반대·혼합), 익명 설문에 답합니다. 중요한 건 — LLM은 "
         "0~100 점수를 직접 매기지 않습니다. 설문 선택지만 고르고 숫자 변환은 코드가 "
         "합니다(판단은 LLM, 단위는 코드). 구버전엔 LLM에 점수를 직접 물었는데, 기분이 "
         "전 점수에 번지는 정서 후광 탓에 비대상 노인이 청년 정책에 수혜 70을 주는 버그가 "
         "있었고, 설문 전환으로 해소했습니다. LangGraph 노드 중 interact(다른 시민을 보고 "
         "입장 변화)는 측정해보니 메인 집계엔 도움이 안 돼, 지금은 SNS 채팅 탭 데모에서만 "
         "씁니다 — 덜어낸 흔적을 숨기지 않고 남겨둔 부분입니다. 화면에 뜨는 게이지·리포트 "
         "같은 집계는 이 단계가 아니라 축3에서 나옵니다. 이 탭은 데모로 직접 보여드립니다.")


def s08_axis_b(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 9, "02 시스템", "축2 · 결과 — 정책 인생극장", accent=ORANGE)
    fx, fw = 0.6, 5.6
    flow_box(slide, fx, 1.8, fw, 0.86, "① 전원 시간경과 시뮬 (1 · 3 · 6개월)",
             "시점마다: 장소 + 경로 + 행동 서사 + 상태 + 막힌 지점",
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=12.5, bsize=10.5)
    arrow_text(slide, fx + fw / 2 - 0.6, 2.68, 1.2, "↓", size=14)
    flow_box(slide, fx, 3.02, fw, 0.86, "② 실제 결과에서 대조 3명 선별",
             "받음 → 수혜 / 진행 중 → 경계 / 막힘·못 닿음 → 사각",
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=12.5, bsize=10.5)
    arrow_text(slide, fx + fw / 2 - 0.6, 3.9, 1.2, "↓", size=14)
    flow_box(slide, fx, 4.24, fw, 0.95, "③ 접근 여정 카드",
             "🧭 알게됨 → 신청 → ⛔막힘  (경로·막힌 지점·삶의 변화)",
             fill=PALE_ORANGE, line_c=ORANGE, tcolor=ORANGE, tsize=12.5, bsize=10.5)
    tf = tb(slide, fx, 5.42, fw, 1.4)
    para(tf, [("채널 5곳이 사각지대의 표현 장치  ",
               {"bold": True, "color": NAVY, "size": 13})], first=True, after=4)
    para(tf, "복지로(온라인) · 주민센터 · 복지관 · 직장/시장 · 집", size=12,
         color=INK, after=3)
    para(tf, "청년은 복지로에서 5분, 고령은 어느 채널에도 닿지 못한다.",
         size=12, bold=True, color=RED, after=0)
    para(tb(slide, fx, 6.5, fw, 0.45),
         "캐시 — 페르소나·시뮬 결과를 재사용 (재호출·재계산 최소화)",
         size=10.5, color=SUB, first=True, after=0)
    shot(slide, 6.5, 1.8, 6.23, 5.0, "‘정책 인생극장’ 탭 캡처",
         "대조 3명 카드 + 여정 띠(알게됨→신청→막힘)가 보이는 화면 추천")
    note(slide, "축 B 정책 인생극장입니다. 시민 전원을 1·3·6개월 시간 경과 속에서 "
         "살게 하고, 그 '실제 결과'에서 수혜·경계·사각 세 사람을 뽑아 카드로 보여줍니다. "
         "카드를 펼치면 '알게 됨 → 신청 → 서류에서 막힘' 같은 여정 띠와 함께, 어떤 "
         "경로로 알게 됐고 어디서 막혔는지가 서사로 나옵니다. 장소 5곳이 사각지대의 "
         "표현 장치인데, 청년은 복지로에서 5분이면 끝나는 일이 고령 어르신에겐 어느 "
         "채널로도 닿지 못하는 일이 됩니다.")


def s08b_axis3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 10, "02 시스템", "축3 · 요약 — 집계", accent=NAVY)
    fx, fw = 0.6, 5.8
    para(tb(slide, fx, 1.75, fw, 0.6),
         "axis3가 축1·축2를 다시 읽어 집계한다 (순수 함수 · LLM 무관)",
         size=15, bold=True, color=NAVY, first=True, after=0)
    flow_box(slide, fx, 2.5, fw, 1.35, "게이지 3종",
             ["정책수용도 · 신청의향지수 · 사회혼란도(=불만 평균)"],
             fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY, tsize=14, bsize=12)
    flow_box(slide, fx, 4.0, fw, 1.35, "종합 리포트 (.md)",
             ["숫자·인용은 코드가 채우고, 산문 4칸만 LLM (1콜)"],
             fill=PALE_NAVY, line_c=NAVY, tcolor=NAVY, tsize=14, bsize=12)
    box(slide, fx, 5.55, fw, 0.85, fill=PALE_GREEN, line=GREEN)
    para(tb(slide, fx + 0.22, 5.67, fw - 0.44, 0.65),
         "화면에 뜨는 모든 숫자는 여기서 — aggregate 노드가 아니라 축3",
         size=12.5, bold=True, color=GREEN, first=True, after=0)
    para(tb(slide, fx, 6.5, fw, 0.45),
         "순수 함수 — 같은 입력이면 재계산 없이 동일 결과 (LLM·캐시 무관)",
         size=10.5, color=SUB, first=True, after=0)
    shot(slide, 6.7, 1.75, 6.03, 5.05, "게이지 · 리포트 화면 캡처",
         "게이지 3종 + 종합 리포트가 보이는 화면")
    note(slide, "축3는 집계 단계입니다. 핵심은 — 화면에 뜨는 게이지와 리포트의 모든 "
         "숫자가 여기 한 곳(axis3)에서 나온다는 점입니다. 원래 LangGraph의 aggregate "
         "노드도 게이지를 계산했지만, 지금은 axis3가 그 값을 덮어쓰기 때문에 실제 화면에 "
         "쓰이는 건 axis3 산출입니다. axis3는 축1의 시민 반응(t0)과 축2의 인생극장 결과를 "
         "입력으로 받아, 순수 파이썬 함수로 집계합니다 — LLM도 streamlit도 타지 않아서 "
         "같은 입력이면 늘 같은 숫자가 나옵니다. 게이지는 정책수용도·신청의향지수·"
         "사회혼란도 셋이고, 사회혼란도는 세 번의 재정의 끝에 '반발 강도 = 불만 점수 "
         "평균'으로 확정했습니다. 종합 리포트(.md)는 숫자와 인용을 코드가 채우고 LLM은 "
         "산문 네 칸만 쓰는 고정 양식이라, 매번 같은 틀로 나옵니다. 집계를 한 곳으로 "
         "모은 게 '판정 지점이 하나라 모순이 안 생긴다'는 재설계의 핵심입니다.")


def s09_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 11, "03 시행착오", "스코프 — 처음 구상에서 최종 형태로",
           accent=ORANGE)
    # 좌: 처음 구상 (흐림, 많아 보이게)
    box(slide, 0.6, 1.8, 5.4, 3.65, fill=GRAY_BG, line=GRAY_BORDER, lw=1.0)
    para(tb(slide, 0.8, 1.95, 5.0, 0.5),
         "처음 구상 — 욕심껏 8가지", size=14, bold=True, color=SUB,
         first=True, after=0)
    feats = ["RAG 게시판", "시민 Agent", "SNS 전파 채팅", "감성 대시보드",
             "전파 그래프", "게시판 자동답변", "A/B 정책 테스트", "쉬운 글 변환"]
    for i, f in enumerate(feats):
        col, row = i % 2, i // 2
        chip(slide, 0.85 + col * 2.55, 2.62 + row * 0.64, f,
             fill=WHITE, color=SUB, w=2.4, h=0.48, size=10.5)
    # 가운데 화살표
    arrow_text(slide, 6.05, 3.0, 1.25, "→", size=30)
    tfa = tb(slide, 5.95, 3.72, 1.45, 0.9)
    para(tfa, "4일 동안", size=10, color=SUB, align=PP_ALIGN.CENTER,
         first=True, after=1)
    para(tfa, "측정하며 좁힘", size=10.5, bold=True, color=ORANGE,
         align=PP_ALIGN.CENTER, after=0)
    # 우: 최종 (선명)
    box(slide, 7.3, 1.8, 5.43, 3.65, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tfr = tb(slide, 7.55, 1.95, 4.95, 3.35)
    para(tfr, "최종 — 검증 가능한 핵심", size=14, bold=True, color=GREEN,
         first=True, after=10)
    para(tfr, [("3축 파이프라인", {"bold": True, "color": NAVY, "size": 13.5})],
         after=2)
    para(tfr, "축1 정보 · 축2 결과 · 축3 요약", size=12, color=INK,
         after=10, lh=1.2)
    para(tfr, [("검증 골격", {"bold": True, "color": NAVY, "size": 13.5})],
         after=2)
    para(tfr, "재료 · 작동 · 한계 (04장)", size=12, color=INK, after=10, lh=1.2)
    para(tfr, [("포지셔닝", {"bold": True, "color": NAVY, "size": 13.5})],
         after=2)
    para(tfr, "‘예측’이 아니라 영향 시나리오", size=12, color=INK, after=0, lh=1.2)
    # 하단 띠
    box(slide, 0.6, 5.7, 12.13, 0.92, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tfb = tb(slide, 0.95, 5.82, 11.5, 0.65, anchor=MSO_ANCHOR.MIDDLE)
    para(tfb, [("좁힌 건 버린 게 아니라 확장 경로 — ",
                {"bold": True, "color": NAVY}),
               ("전파는 미리마을로, RAG는 게시판으로 결국 돌아왔다 (05장)",
                {"color": INK})], size=13.5, first=True, after=0,
         align=PP_ALIGN.CENTER)
    note(slide, "미리랩은 한 번에 설계되지 않았습니다. 처음엔 욕심껏 여덟 가지 기능을 "
         "구상했지만 — RAG 게시판, 시민 에이전트, SNS 전파, 감성 대시보드, 전파 그래프, "
         "게시판 자동답변, A/B 정책 테스트, 쉬운 글 변환 — 4일 안에 '검증 가능한 핵심'만 "
         "남기기로 했습니다. 왼쪽이 처음 욕심, 오른쪽이 최종 형태입니다. 좁히는 과정엔 "
         "분명한 판단이 있었습니다. 전파 그래프는 통계에서 무작위로 뽑은 페르소나가 서로 "
         "모르는 독립 개인이라 그 사이에 전파선을 그릴 근거가 없어서 버렸고, '예측'이라는 "
         "단어는 '맞히는 도구냐'는 검증 공격을 부르기 때문에 '영향 시나리오'로 포지셔닝을 "
         "바꿨습니다. 중요한 건 — 좁힌 것들이 버려진 게 아니라 뒤에서 다른 모습으로 "
         "돌아온다는 점입니다. 전파는 미리마을로, RAG는 게시판으로 돌아왔습니다(05장). "
         "다음은 이렇게 좁히는 과정에서 만난 구체적인 시행착오 세 가지입니다.")


def s10_two_judges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 12, "03 시행착오", "시행착오 ① — 두 심판 문제",
           accent=ORANGE)
    tf = tb(slide, 0.6, 1.8, 12.13, 1.5)
    para(tf, [("증상  ", {"bold": True, "color": RED, "size": 14}),
              ("‘경계’로 뽑힌 신민재의 카드 서사가 “월세를 받았다” — 라벨과 이야기가 모순",
               {"size": 14})], first=True, after=6)
    para(tf, [("원인  ", {"bold": True, "color": RED, "size": 14}),
              ("심판이 둘 — 라벨은 시뮬 ", {"size": 14}),
              ("전", {"bold": True, "size": 14}),
              (" 점수 매트릭스가, 이야기는 시뮬 ", {"size": 14}),
              ("후", {"bold": True, "size": 14}),
              (" LLM이 정한다. 두 심판은 서로 어긋난다.", {"size": 14})], after=6)
    para(tf, [("결론  ", {"bold": True, "color": BLUE, "size": 14}),
              ("수혜/경계/사각은 점수가 아니라 ", {"size": 14}),
              ("결과(받음 / 막힘 / 못 닿음)", {"bold": True, "color": BLUE, "size": 14}),
              ("다.", {"size": 14})], after=0)
    # Before / After
    box(slide, 0.6, 3.55, 5.9, 2.9, fill=PALE_RED, line=RED, lw=1.0)
    tf1 = tb(slide, 0.85, 3.75, 5.4, 2.5)
    para(tf1, "Before — 수치 예측 선별", size=14.5, bold=True, color=RED,
         first=True, after=6)
    bullets(tf1, [
        ("시뮬 전에 점수 매트릭스로 3명을 미리 고른다", {"bullet_color": RED}),
        ("라벨 따로, 서사 따로 → 구조적으로 어긋남", {"bullet_color": RED}),
        ("반응 점수로 뽑아도 같은 함정 — 임수빈은 자기 수혜성을 18점으로 "
         "매기지만 실제론 요건에 막힌 사각지대", {"bullet_color": RED}),
    ], size=12, after=5, first_done=True)
    box(slide, 6.83, 3.55, 5.9, 2.9, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf2 = tb(slide, 7.08, 3.75, 5.4, 2.5)
    para(tf2, "After — 결과 기반 선별", size=14.5, bold=True, color=GREEN,
         first=True, after=6)
    bullets(tf2, [
        ([("전원을 먼저 살게 한 뒤", {"bold": True}),
          (", 실제 궤적의 결과에서 3명을 고른다", {})], {"bullet_color": GREEN}),
        ("라벨과 이야기가 같은 궤적에서 나옴 → 어긋날 수 없음", {"bullet_color": GREEN}),
        ("심판이 하나로 통일됨. 전원 시뮬 비용도 확인 결과 문제없음", {"bullet_color": GREEN}),
    ], size=12, after=5, first_done=True)
    note(slide, "첫 번째 시행착오, 저희가 '두 심판 문제'라고 부르는 버그입니다. 처음엔 "
         "시뮬레이션 전에 점수로 세 명을 '예측'해서 뽑았는데, 경계로 뽑힌 사람의 "
         "이야기가 '월세를 받았다'로 나오는 모순이 생겼습니다. 라벨을 정하는 심판과 "
         "이야기를 쓰는 심판이 달랐던 거죠. 핵심 통찰은 — 수혜·경계·사각은 점수가 "
         "아니라 '결과'라는 것. 그래서 전원을 먼저 살게 한 뒤 실제 결과에서 뽑는 "
         "방식으로 바꿨고, 이제 라벨과 이야기는 구조적으로 어긋날 수 없습니다.")


def s11_billion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 13, "03 시행착오", "시행착오 ② — “전 국민 10억 지급” 버그",
           accent=ORANGE)
    steps = [
        ("실험", "극단 정책 “전 국민에게 10억 지급” 입력", RED, PALE_RED),
        ("증상", "전원 찬성 · 혼란도 21 — 시민이 ‘이해득실 계산기’처럼 내 혜택만 따짐",
         RED, PALE_RED),
        ("처방", "프롬프트 최상단에 “현실적으로 판단하라” — 재원·물가·증세·형평성·"
         "과대공약 의심 (+ 합리적 정책엔 냉소 금지)", BLUE, PALE_NAVY),
        ("결과", "10억 지급 → 전원 반대로 전환. 현실성 판단 작동 확인", GREEN, PALE_GREEN),
    ]
    y = 1.8
    for tag, body, c, pale in steps:
        chip(slide, 0.6, y + 0.12, tag, fill=c, w=1.1, h=0.34)
        box(slide, 1.9, y, 10.83, 0.62, fill=pale, line=c, lw=0.75)
        tf = tb(slide, 2.1, y + 0.05, 10.4, 0.52, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, body, size=12.5, color=INK, first=True, after=0)
        y += 0.78
    box(slide, 0.6, 5.0, 12.13, 1.7, fill=PALE_ORANGE, line=ORANGE, lw=1.0)
    tf = tb(slide, 0.95, 5.2, 11.5, 1.35)
    para(tf, [("그 후 — 이 처방은 결국 철거됐다", {"bold": True, "color": ORANGE,
                                                "size": 14.5})], first=True, after=5)
    para(tf, "‘현실적으로 판단하라’를 넣자 정상 정책의 혼란도까지 오르는 과교정 발생 → "
         "프롬프트 전면 재설계(판단 지시 0줄, 무대만 제공)로 지시 자체를 철거.",
         size=13, after=3, lh=1.25)
    para(tf, [("답은 지시가 아니라 측정(설문 전환·자가판정)과 모델 선택(Gemini)이었다. ",
               {"size": 13}),
              ("교훈: 프롬프트로 조향하지 마라.", {"bold": True, "size": 13}),
              ("  분포는 데이터와 측정으로 푼다.", {"size": 13, "color": SUB})], after=0)
    note(slide, "두 번째 시행착오 — 이건 두 막짜리입니다. 1막: 극단 정책 '전 국민 10억 "
         "지급'에 전원 찬성이 나왔습니다. 시민들이 '내가 받느냐'만 따지는 계산기였던 "
         "거죠. 프롬프트에 '현실적으로 판단하라'를 넣자 전원 반대로 바뀌어 해결된 듯 "
         "보였습니다. 2막: 그 처방이 정상 정책의 혼란도까지 끌어올리는 과교정을 만들었고, "
         "결국 프롬프트를 전면 재설계하면서 판단 지시를 전부 철거했습니다 — 프롬프트는 "
         "무대만 제공하고, 판단은 페르소나 데이터에서 나오게. 진짜 답은 지시가 아니라 "
         "측정 방식(점수를 익명 설문로)과 모델 선택(Gemini 전환)이었습니다. '프롬프트로 "
         "조향하지 마라'가 이 시행착오의 최종 교훈이고, 이 원칙 위에서 다음 장의 검증이 "
         "돌아갑니다.")


def s12_guards(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 14, "03 시행착오", "시행착오 ③ — 판단은 LLM에게, 불변식은 코드로",
           accent=ORANGE)
    box(slide, 0.6, 1.75, 12.13, 0.85, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tf = tb(slide, 0.95, 1.9, 11.5, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("LLM은 확률적이다 — 서사와 판단은 LLM이, ", {"size": 15}),
              ("어겨선 안 되는 규칙은 코드 가드가 보증한다", {"bold": True, "color": NAVY, "size": 15})],
         first=True, after=0, align=PP_ALIGN.CENTER)
    cards = [
        ("상태 비퇴행 가드", [
            "버그: ‘신청 → 모름’으로 궤적이 퇴행 (서류에 막혀 포기한 걸 LLM이 ‘몰랐다’로 오기)",
            "가드: 상태는 되돌아가지 않는다 — applied→unaware면 ‘막힘’으로 자동 교정",
            "결과: 조유정이 ‘자격은 되는데 서류에 막힘’ 사각 사례로 바르게 남음"],
         GREEN, PALE_GREEN),
        ("태그 대상 게이트", [
            "버그: 46세 박성인이 ‘청년’ 정책의 사각지대로 오라벨",
            "가드: 대상 여부 = 태그(나이·소득·가구)로 판정하는 ‘사실’, 결과 = LLM 궤적 — 둘을 분리",
            "결과: 비대상자는 ‘무관’으로 분류"],
         ORANGE, PALE_ORANGE),
        ("전파 게이트", [
            "무대: 잠시 후 소개할 ‘미리마을’ — 정책 소식이 사람을 타고 퍼지는 시뮬",
            "위험: LLM이 근거 없이 ‘어느새 다 알게 됨’으로 부풀릴 수 있음",
            "가드: 모르던 사람은 ‘아는 사람과 만났을 때만’ 새로 안다 — 코드가 강제",
            "결과: 전파 경로 추적이 신뢰 가능해짐 (인지 누수 0)"],
         RED, PALE_RED),
    ]
    x = 0.6
    for title, lines, c, pale in cards:
        box(slide, x, 2.85, 3.93, 3.75, fill=pale, line=c, lw=1.0)
        tf = tb(slide, x + 0.2, 3.05, 3.55, 3.4)
        para(tf, title, size=14, bold=True, color=c, first=True, after=7)
        for i, ln in enumerate(lines):
            tag = ["", "", ""][i] if False else None
            para(tf, ln, size=10.8, color=INK, after=7, lh=1.22)
        x += 4.1
    note(slide, "세 번째는 버그 하나가 아니라, 반복해서 발견한 설계 패턴입니다. LLM은 "
         "확률적이라 가끔 말이 안 되는 출력을 냅니다 — 신청했던 사람이 '몰랐다'로 "
         "퇴행하거나, 46세를 청년 정책 사각지대로 잘못 라벨하거나, 모르던 사람이 "
         "어느새 다 알게 되거나. 저희 답은 일관됩니다: 서사와 판단은 LLM에게 맡기되, "
         "어겨선 안 되는 규칙(상태는 안 되돌아간다 / 대상 여부는 사실이다 / 모르는 "
         "사람은 만나야 안다)은 코드 가드가 보증합니다. 이 철학이 시스템 전체를 "
         "관통합니다.")


def s13_validation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 15, "04 검증", "검증 — ‘맞혔나’가 아니라 ‘합리적으로 작동하는가’",
           accent=GREEN, title_size=22)
    box(slide, 0.6, 1.62, 12.13, 1.16, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tfb = tb(slide, 0.95, 1.72, 11.5, 0.98, anchor=MSO_ANCHOR.MIDDLE)
    para(tfb, [("미리랩은 미래를 맞히는 예측기가 아니라, 정책의 파장을 미리 보는 시뮬레이터다. ",
                {"bold": True, "color": NAVY, "size": 13}),
               ("‘여론조사를 얼마나 맞혔나’로 평가하지 않는다 — 그 지표는 약하고, 솔직히 인정한다.",
                {"size": 12.5, "color": INK})], first=True, after=3, lh=1.25)
    para(tfb, "대신 — 말도 안 되는 극단 정책에도 시민이 상식적으로 반응하면, 일반 정책에서의 반응도 믿을 수 있다.",
         size=12.5, bold=True, color=GREEN, after=0, lh=1.25)
    rows = [
        ("① 재료", "시민 24명이 한국 인구를 닮았나?", "분포 7/7 통과 — 무작위 추출 "
         "기대범위 안 (부트스트랩 1만회, LLM 0콜)", GREEN, PALE_GREEN),
        ("② 작동", "반응 차이를 만드는 게 페르소나인가?", "사전등록 0/4 — 원인은 "
         "채점 기준이었다 (다음 장)", ORANGE, PALE_ORANGE),
        ("③ 강건성", "말도 안 되는 정책에도 상식을 지키나?", "행동 벤치마크 7/7 통과 — "
         "극단 정책 5종에서 상식 방향으로 (다다음 장)", GREEN, PALE_GREEN),
    ]
    x = 0.6
    for tag, q, res, c, pale in rows:
        box(slide, x, 2.92, 3.93, 1.5, fill=pale, line=c, lw=1.0)
        tf = tb(slide, x + 0.2, 3.06, 3.55, 1.25)
        para(tf, [(tag + "  ", {"bold": True, "color": c, "size": 13.5}),
                  (q, {"size": 11.5, "color": INK})], first=True, after=4, lh=1.2)
        para(tf, res, size=11, color=INK, after=0, lh=1.22)
        x += 4.1
    img_fit(slide, ROOT / "eval" / "persona_eval_viz.png", 0.6, 4.5, 12.13, 2.45)
    note(slide, "이제 '이걸 믿을 수 있나'에 답합니다. 먼저 분명히 — 미리랩은 미래를 "
         "맞히는 예측기가 아닙니다. 정책을 넣으면 어떤 일이 벌어질 수 있는지, 그 "
         "가능성과 서사를 보는 시뮬레이터입니다. 그래서 '과거 여론조사를 얼마나 맞혔나' "
         "같은 정확도로는 평가하지 않습니다 — 그 지표는 약할 수밖에 없고, 솔직히 "
         "인정합니다. 대신 저희가 기댄 논리는 스트레스 테스트입니다. 말도 안 되는 극단 "
         "정책을 넣어도 시민들이 상식적인 방향으로 반응한다면, 평범한 정책에서의 반응도 "
         "믿을 수 있다는 겁니다. 그 위에서 세 가지를 봅니다 — 재료(시민이 한국 인구를 "
         "닮았는가), 작동(반응 차이를 만드는 게 페르소나인가), 강건성(극단 입력에도 "
         "상식을 지키는가). 방법론은 사전등록을 따릅니다 — 판정 기준을 실행 전에 박아두고, "
         "결과가 나빠도 점수를 고치지 않습니다. 아래 그림이 ① 재료 검증입니다. 전체 "
         "100만 명 분포(회색)와 우리 24명(파랑)을 변수별로 비교해, 무작위로 뽑아도 생기는 "
         "어긋남의 정상 범위인지 부트스트랩 1만 번으로 판정했고, 7개 변수 전부 "
         "통과했습니다. 데이터셋 자체의 현실 정합은 NVIDIA가 카드에 검증·문서화한 것을 "
         "인용합니다.")


def s14_prereg(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 16, "04 검증", "검증 ② — 사전등록 0/4, 원인은 채점 기준이었다",
           accent=GREEN)
    tf = tb(slide, 0.6, 1.7, 5.85, 2.55)
    para(tf, [("실험  ", {"bold": True, "color": NAVY, "size": 13}),
              ("같은 청년월세 정책을 인물 카드 ON / OFF(익명)로 — "
               "카드가 있을 때만 답이 자기 처지를 따라가야 함", {"size": 12.5})],
         first=True, after=6, lh=1.25)
    para(tf, [("결과  ", {"bold": True, "color": RED, "size": 13}),
              ("사전등록 점수 0/4 탈락 — 그대로 기록", {"size": 12.5, "bold": True})],
         after=6)
    para(tf, [("분석  ", {"bold": True, "color": BLUE, "size": 13}),
              ("우리 채점 기준(나이·소득)이 ‘대상’이라 한 10명 전원이 자기를 비대상으로 "
               "정확 판정 — 원문의 숨은 요건(별도거주·무주택·소득기준)까지 읽었다. "
               "실제 청년월세도 신청자 2/3가 요건 탈락.", {"size": 12.5})],
         after=0, lh=1.28)
    # 우측 시민 인용
    box(slide, 6.7, 1.7, 6.03, 2.55, fill=CARD, line=LINE)
    tf2 = tb(slide, 6.95, 1.86, 5.55, 2.25)
    para(tf2, "시민들이 직접 말한 탈락 사유", size=13, bold=True, color=NAVY,
         first=True, after=6)
    para(tf2, [("천명준(19)  ", {"bold": True, "color": BLUE, "size": 11.5}),
               ("“부모님이랑 따로 살아야 주는 거잖아요. 저는 해당이 안 될 텐데…”",
                {"size": 11.5})], after=5, lh=1.2)
    para(tf2, [("심석현(26)  ", {"bold": True, "color": BLUE, "size": 11.5}),
               ("“기사 끝에 보니까 주택 소유자는 제외라고 딱 적혀 있어요.”",
                {"size": 11.5})], after=5, lh=1.2)
    para(tf2, [("조유정(31)  ", {"bold": True, "color": BLUE, "size": 11.5}),
               ("“중위소득 60%면 생각보다 낮거든요. 저는 기준에서 걸려요.”",
                {"size": 11.5})], after=0, lh=1.2)
    img_fit(slide, ROOT / "eval" / "ablation_shift_viz.png", 0.6, 4.35, 12.13, 2.6)
    note(slide, "검증 ②의 이야기는 저희 발표에서 가장 정직한 대목입니다. 기준을 먼저 "
         "등록하고 돌렸더니 0대 4로 떨어졌습니다 — 그리고 점수는 그대로 둔 채 원자료를 "
         "부검했더니, 떨어진 이유가 반전이었습니다. 우리 채점 기준은 나이와 소득만 보고 "
         "젊은 10명을 '대상자'라 표시했는데, 시뮬 속 시민들은 정책 원문의 별도거주·"
         "무주택·소득 요건까지 읽고 전원이 '나는 해당 안 됨'을 정확히 판정한 겁니다. "
         "실제 청년월세도 신청자 3명 중 2명이 요건에서 탈락했으니, 시뮬이 현실을 맞게 "
         "재현했고 채점표가 그걸 몰랐던 거죠. 아래 그림이 진짜 증거입니다 — 익명이면 "
         "'나도 도움되겠지'로 낙관하던 응답(회색)이, 인물 카드가 들어가는 순간 자기 "
         "처지(영향 없음 99%)로 수렴합니다. 페르소나가 차이를 만든다는 건 분산이 아니라 "
         "이 '평균 이동'으로 증명됐습니다. 사전등록 원칙대로 점수는 0/4 그대로 두고, "
         "채점 도구 재설계는 다음 사이클로 남겼습니다 — 기준이 틀렸다고 기준을 몰래 "
         "고치면 검증이 아니니까요.")


def s14b_bench(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 17, "04 검증", "검증 ③ — 행동 벤치마크: 가상 정책 5종, 7/7 통과",
           accent=GREEN)
    tf = tb(slide, 0.6, 1.66, 5.85, 2.6)
    para(tf, [("가상 정책 5종 (일부러 부조리하게)  ", {"bold": True, "color": NAVY,
                                                  "size": 13})], first=True, after=4)
    bullets(tf, [
        ("주 3.5일 근무제 — 직업 축 / 반려묘 전 국민 보급 — 순효용 축", {}),
        ("경로 무임승차 폐지 — 세대 축 / 남 +10만·여 −10만 — 성별 축", {}),
        ("고소득자 의무 기부 — 소득 축 (시민이 자기 형편을 자가판정)", {}),
    ], size=11.8, after=4, first_done=True)
    para(tf, [("사전등록 체크 7/7 통과", {"bold": True, "color": GREEN, "size": 14}),
              ("  — 갈등이 매번 상식적인 집단 축으로 갈라짐", {"size": 12})],
         after=6)
    para(tf, [("사례 ①  ", {"bold": True, "color": ORANGE, "size": 12.5}),
              ("반려묘 유일 찬성자 = 페르소나에 ‘펫푸드 창업 꿈’이 있는 신민재 — "
               "입장이 인구통계가 아니라 서사와 정렬", {"size": 12})], after=4, lh=1.25)
    para(tf, [("사례 ②  ", {"bold": True, "color": ORANGE, "size": 12.5}),
              ("남성에게 돈을 주는 정책에 남성 23%가 반대 — “내가 받고 아내가 내면 "
               "우리 집은 0원” (가구 단위 공정성)", {"size": 12})], after=0, lh=1.25)
    img_fit(slide, ROOT / "eval" / "behavior_bench_viz.png", 6.6, 1.75, 6.13, 4.0)
    box(slide, 6.6, 5.9, 6.13, 0.85, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf3 = tb(slide, 6.8, 6.0, 5.75, 0.65, anchor=MSO_ANCHOR.MIDDLE)
    para(tf3, "이 5종은 데모 녹화본 — 발표 자리에서 키·네트워크 없이 그대로 재생 시연 가능",
         size=11.5, bold=True, color=GREEN, first=True, after=0,
         align=PP_ALIGN.CENTER)
    note(slide, "검증 ③은 행동 벤치마크 — 일부러 만든 가상 정책 5종(일부는 농담조)을 "
         "넣고 사회가 상식 방향으로 갈라지는지 보는 고정 시험지입니다. 기대 방향은 "
         "실행 전에 등록했고, 7개 체크 전부 통과했습니다. 무임승차 폐지엔 60세 이상이 "
         "반대 57%로 갈라지고(미만은 6%), 여성에게서 걷어 남성에게 주는 부조리 정책엔 "
         "여성 전원이 반대합니다. 백미 두 개 — 전 국민 반려묘 보급의 유일한 찬성자는 "
         "페르소나 서사에 펫푸드 창업 꿈이 있는 시민이었고, 남성에게 돈을 주는 정책에 "
         "남성의 23%가 '우리 집 전체로는 0원'이라며 반대했습니다. 입장이 인구통계를 "
         "넘어 서사·가구 단위 사고와 정렬된다는 증거입니다. 그리고 이 벤치마크는 "
         "일회용이 아니라 — 모델이나 프롬프트를 바꿀 때마다 다시 돌려 '여전히 사람답게 "
         "구는가'를 확인하는 회귀 테스트로 씁니다. 5종 전부 데모 녹화본이라 지금 이 "
         "자리에서 키 없이 재생 시연도 가능합니다.")


def s15_board(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 18, "05 확장", "게시판 RAG — 문서 근거로 답하는 정책 문의 게시판",
           accent=RED)
    tf = tb(slide, 0.6, 1.8, 5.7, 4.9)
    para(tf, "정책 문의 게시판 + 문서 근거 답변", size=15.5, bold=True,
         color=NAVY, first=True, after=8)
    bullets(tf, [
        ([("질문 → 검색된 근거 ‘안에서만’ 답변", {"bold": True}),
          ("  (환각 억제)", {"color": SUB})], {}),
        ("현재 정책은 항상 기본 근거 + PDF/TXT/MD 업로드로 문서 확장", {}),
        ("Chroma 벡터 검색 + OpenAI 임베딩 → 근거·품질지표 함께 표시", {}),
        ("API 키가 없으면 추출식 폴백 — 데모는 항상 동작", {}),
    ], size=12.5, after=7, first_done=True)
    para(tf, "", size=4, after=2)
    box(slide, 0.6, 4.7, 5.7, 1.95, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf2 = tb(slide, 0.85, 4.88, 5.2, 1.6)
    para(tf2, "협업 방식 — 자리만 비워두고 끼워 넣기", size=13.5, bold=True, color=GREEN,
         first=True, after=5)
    para(tf2, "본체에는 끼울 자리와 반환 계약만 미리 설계해 두고, 팀원이 "
         "독립 패키지로 RAG 엔진을 구현해 마지막에 끼워 넣었다. 모듈 경계를 "
         "지킨 덕에 본체는 한 줄도 깨지지 않았다.", size=11.8, color=INK, after=0, lh=1.3)
    shot(slide, 6.6, 1.8, 6.13, 4.85, "‘게시판’ 탭 캡처",
         "질문 → 답변 + 📎근거 + 품질지표 표가 보이는 화면 추천")
    note(slide, "확장 파트 첫 번째, 게시판 RAG입니다. 2일차에 'MVP에서 제외'로 좁혔던 "
         "RAG가 마지막에 팀의 손으로 돌아왔습니다. 시민이나 입안자가 정책 문서에 "
         "질문하면, 벡터 검색으로 찾은 근거 '안에서만' 답하고 근거와 품질지표를 함께 "
         "보여줍니다. 협업 방식이 포인트인데 — 본체에는 자리와 반환 계약만 설계해 두고 "
         "팀원이 독립 패키지로 구현해서 마지막에 끼워 넣었습니다. 키가 없으면 추출식 "
         "폴백으로 동작하니 데모도 안전합니다.")


def s16_village1(prs, sprites):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 19, "05 확장", "미리마을 — 시민이 살아가는 마을 (Generative Agents)",
           accent=RED)
    tf = tb(slide, 0.6, 1.75, 6.1, 3.1)
    para(tf, [("Park et al. (2023) ", {"bold": True, "color": BLUE}),
              ("의 원형 — 에이전트가 공간에서 자기 스케줄대로 살아가는 마을",
               {})], size=13.5, first=True, after=8, lh=1.25)
    bullets(tf, [
        ([("10명 손제작 고정 캐스트", {"bold": True}),
          (" — 무작위 통계 샘플엔 ‘관계’가 없다. 전파 무대엔 서로 아는 마을이 필요 "
           "(축 A·B와 목적이 다른 의도된 선택)", {})], {}),
        ("LLM이 각자의 하루 일과(시간 → 장소 → 행동)를 생성", {}),
        ([("녹화 후 재생 — 하루를 미리 생성, 브라우저는 LLM 0콜로 재생",
           {"bold": True}), ("  (즉시·반복·오프라인 = 발표 안전망)", {"color": SUB})], {}),
        ("도로 waypoint 652개 → 그래프 + BFS 길찾기로 실제 ‘도보’ 이동", {}),
        ("마주치면 — 그날 동선·관계가 묻어나는 맥락 대화", {}),
    ], size=12, after=5.5, first_done=True)
    box(slide, 0.6, 4.95, 6.1, 1.7, fill=PALE_ORANGE, line=ORANGE, lw=1.0)
    tf2 = tb(slide, 0.85, 5.1, 5.6, 1.45)
    para(tf2, "시행착오 — 스케줄은 순간이동, 재생은 도보", size=13,
         bold=True, color=ORANGE, first=True, after=4)
    para(tf2, "LLM이 짠 스케줄은 이동시간 0을 가정 → 막상 재생하니 만남 14개 중 "
         "7개가 ‘도착 전에 끝남’. 실제 걸어서 도착하는 시간을 시뮬해 만남을 "
         "다시 뽑자 발화 6/6 보장.", size=11.5, color=INK, after=0, lh=1.28)
    # 우측: 마을 맵 + 창발 한 줄
    img_fit(slide, ROOT / "미리마을" / "assets" / "map.png", 6.95, 1.75, 5.8, 4.0)
    box(slide, 6.95, 5.85, 5.8, 0.8, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf3 = tb(slide, 7.15, 5.95, 5.4, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    para(tf3, [("창발 —  ", {"bold": True, "color": GREEN, "size": 12}),
               ("시키지 않았는데 — 민수·수아가 광장 산책을 약속하고, 같은 카페 "
                "사장이 세 손님과 ‘각자의 하루’로 다른 대화를 한다",
                {"size": 11.3})], first=True, after=0, lh=1.2)
    note(slide, "마지막 확장, 미리마을입니다. Generative Agents 논문의 원형 — "
         "에이전트가 공간에서 스케줄대로 살아가는 마을을 포켓몬식 오버월드로 "
         "구현했습니다. LLM이 10명 각자의 하루를 생성하면 브라우저는 LLM 호출 없이 "
         "재생만 합니다. 발표 중에도 안전하죠. 재미있는 시행착오가 있었는데, LLM이 짠 "
         "스케줄은 순간이동을 가정해서 막상 걸어가면 만남의 절반이 증발했습니다. 실제 "
         "도보 시간을 시뮬해서 만남을 다시 뽑았고요. 그리고 시키지 않았는데 캐릭터들이 "
         "서로를 일정에 엮고, 같은 카페 사장이 손님마다 다른 대화를 하는 '창발'이 "
         "나타났습니다. [Q&A 대비 — \"손으로 안 만든다더니 모순 아니냐\": 축 A·B의 "
         "통계 샘플은 '대표성'이 목적이라 무작위 추출이 맞고, 미리마을은 '전파'가 "
         "목적이라 서로 아는 관계가 필수입니다. 무작위 샘플엔 관계가 없으니 손제작 — "
         "목적이 다르면 재료도 다릅니다.]")


def s17_village2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 20, "05 확장", "미리마을 × 정책 — 소식은 누구에게, 어떤 경로로 닿는가",
           accent=RED)
    box(slide, 0.6, 1.75, 12.13, 0.8, fill=PALE_NAVY, line=NAVY, lw=1.0)
    tf = tb(slide, 0.95, 1.88, 11.5, 0.58, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("질문의 전환:  “정책이 각자에게 어떤 효과인가”  →  ",
               {"size": 14.5}),
              ("“정책 ‘소식’이 누구에게 · 어떤 경로로 닿는가”",
               {"bold": True, "color": NAVY, "size": 14.5})], first=True, after=0,
         align=PP_ALIGN.CENTER)
    fx = 0.6
    flow_box(slide, fx, 2.8, 2.7, 1.5, "① 시드 주입",
             "담당 공무원 영희 1명만 알고 시작", fill=PALE_RED, line_c=RED,
             tcolor=RED, tsize=12.5, bsize=11)
    arrow_text(slide, fx + 2.72, 3.35, 0.5, "→", size=16)
    flow_box(slide, fx + 3.24, 2.8, 2.7, 1.5, "② 만남 전파",
             "스케줄대로 살다 마주친 대화에서 소식이 옮음", fill=PALE_ORANGE,
             line_c=ORANGE, tcolor=ORANGE, tsize=12.5, bsize=11)
    arrow_text(slide, fx + 5.96, 3.35, 0.5, "→", size=16)
    flow_box(slide, fx + 6.48, 2.8, 2.7, 1.5, "③ 밤 일기 (reflection)",
             "하루를 일기로 압축 → 다음날의 기억으로", fill=PALE_NAVY, line_c=BLUE,
             tcolor=BLUE, tsize=12.5, bsize=11)
    arrow_text(slide, fx + 9.2, 3.35, 0.5, "→", size=16)
    flow_box(slide, fx + 9.72, 2.8, 2.41, 1.5, "④ [▶다음날]",
             "며칠이고 이어서 — 전파의 시간 구조", fill=PALE_GREEN, line_c=GREEN,
             tcolor=GREEN, tsize=12.5, bsize=11)
    # 결과 숫자
    box(slide, 0.6, 4.55, 5.9, 2.1, fill=CARD, line=LINE)
    tf2 = tb(slide, 0.85, 4.72, 5.4, 1.8)
    para(tf2, "1일차 실측 (실제 LLM)", size=13.5, bold=True, color=NAVY,
         first=True, after=6)
    bullets(tf2, [
        ([("10명 중 9명에게 전파", {"bold": True}),
          (" — 공무원 시드 → 카페 사랑방 → 놀이터", {"color": SUB})], {}),
        ([("사각지대 1명", {"bold": True, "color": RED}),
          (" — 동선이 겹치지 않은 박어르신", {"color": SUB})], {}),
        ([("근거 없는 인지 누수 0", {"bold": True, "color": GREEN}),
          (" — 전파 게이트가 코드로 보증", {"color": SUB})], {}),
    ], size=12, after=5, first_done=True)
    box(slide, 6.83, 4.55, 5.9, 2.1, fill=PALE_GREEN, line=GREEN, lw=1.0)
    tf3 = tb(slide, 7.08, 4.72, 5.4, 1.8)
    para(tf3, "폐기했던 전파 기능의 부활", size=13.5, bold=True,
         color=GREEN, first=True, after=6)
    para(tf3, "축 A의 전파 그래프는 ‘서로 모르는 통계 샘플’이라 버렸다(03장). "
         "미리마을은 서로 아는 고정 캐스트 — 여기서의 전파는 연출이 아니라 "
         "관계와 동선의 결과다.", size=12, color=INK, after=0, lh=1.3)
    t_demo = tb(slide, 0.6, 6.72, 12.13, 0.36)
    para(t_demo, "🎬 여기서 라이브 데모 — 미리마을 재생 60~90초 (재생은 LLM 0콜 "
         "= 네트워크가 죽어도 안전)", size=12.5, bold=True, color=RED,
         align=PP_ALIGN.CENTER, first=True, after=0)
    note(slide, "미리마을에 정책을 주입하면 질문이 바뀝니다. '효과가 어떤가'가 아니라 "
         "'소식이 누구에게, 어떤 경로로 닿는가'. 주민센터 정책 담당 공무원 영희, 한 "
         "사람만 직무로 아는 상태로 시작하면, 마을 사람들이 스케줄대로 살다가 마주친 "
         "대화에서 소식이 옮습니다. 밤에는 각자 일기로 하루를 압축해 다음날의 기억이 "
         "됩니다. 실측 결과 하루 만에 "
         "10명 중 9명에게 퍼졌는데, 동선이 안 겹친 박어르신 한 분이 사각지대로 남았습니다. "
         "12장에서 보신 '전파 게이트'가 바로 이 무대의 코드 가드입니다. 그리고 중요한 것 — "
         "처음에 버렸던 '전파'가 여기서 부활했습니다. 여긴 서로 아는 마을이라 전파가 "
         "연출이 아니라 관계와 동선의 결과이기 때문입니다. → 이 슬라이드 직후 미리마을 "
         "재생 데모(60~90초): 발표 첫 장에서 약속한 '이 친구들이 직접 움직이는 모습'을 "
         "여기서 회수합니다. 녹화 재생이라 LLM 0콜, 발표장 네트워크와 무관하게 안전.")


def s18_discipline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 21, "06 마무리", "4일 · 4명 — 협업 규율 6가지",
           accent=SUB)
    cards = [
        ("state.py = 팀 계약", "공유 스키마는 기존 필드 변경 금지, 추가만(additive). "
         "4명이 같은 계약을 읽고 쓴다."),
        ("mock 모드 = 발표 안전망", "API 키·네트워크 없이 전체 UI 시연. "
         "발표 중 장애에도 데모는 죽지 않는다."),
        ("import 부작용 0", "어떤 모듈도 import 시 네트워크/LLM 호출 없음. "
         "실행은 버튼을 눌렀을 때만."),
        ("헤드리스 회귀 스크립트", "선별·게이지·히트맵·전체 앱 스모크 — "
         "변경할 때마다 키 없이 자동 검증."),
        ("탭 격리", "한 탭이 죽어도 다른 탭은 산다 — 각 탭을 try/except로 격리."),
        ("XSS 방어", "LLM 자유 텍스트·시민 이름은 화면에 닿기 전 전부 escape."),
    ]
    positions = [(0.6, 1.85), (6.83, 1.85), (0.6, 3.35), (6.83, 3.35),
                 (0.6, 4.85), (6.83, 4.85)]
    for (title, body), (x, y) in zip(cards, positions):
        box(slide, x, y, 5.9, 1.35, fill=CARD, line=LINE)
        tf = tb(slide, x + 0.22, y + 0.13, 5.5, 1.1)
        para(tf, title, size=13.5, bold=True, color=NAVY, first=True, after=3)
        para(tf, body, size=11, color=INK, after=0, lh=1.2)
    t = tb(slide, 0.6, 6.4, 12.13, 0.5)
    para(t, "역할 분담 — (발표 전 채워주세요: 데이터/페르소나 · Agent/프롬프트 · "
         "대시보드/지표 · UI/시각화)", size=11.5, color=SUB,
         align=PP_ALIGN.CENTER, first=True, after=0)
    note(slide, "기능 얘기를 마치기 전에, 4일 동안 4명이 서로의 작업을 깨지 않게 한 "
         "규율을 짧게 소개합니다. 공유 스키마는 변경 금지·추가만 허용하는 팀 계약으로 "
         "묶었고, 키 없이 도는 mock 모드와 헤드리스 회귀 스크립트 덕분에 누가 무엇을 "
         "바꿔도 안전하게 검증했습니다. 게시판 RAG가 마지막 날 본체를 한 줄도 안 깨고 "
         "들어올 수 있었던 게 이 규율의 실증입니다. (시간이 빠듯하면 이 장은 한 문장으로 "
         "줄이고 넘어가도 됩니다. 역할 분담 줄은 실제 팀 구성으로 채워주세요.)")


def s19_limits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 22, "06 마무리", "한계 — 먼저 말해두는 다섯 가지",
           accent=SUB)
    items = [
        ("미래를 맞히는 도구가 아니다", "일관되게 ‘영향 시나리오’로 포지셔닝. "
         "신뢰성은 예측 정확도가 아니라 검증 3축(재료·작동·한계)과 사전등록으로 방어한다."),
        ("소득 판정은 근사다", "직업→소득 버킷은 휴리스틱. 그래서 검증 ⑤(의무 기부)는 "
         "시민의 자가판정을 썼다 — 코드보다 시민이 자기 형편을 더 잘 안다."),
        ("‘실현 기대’ 측정축이 없다", "“안 믿지만 받으면 좋지”가 수혜·의향 점수에 "
         "섞인다 — 설문에 실현 기대 문항 추가가 다음 후보."),
        ("표본이 작다", "시민 24명 — 비율 최소 단위 4.2%p. 통계적 일반화가 아니라 "
         "‘구조가 나타나는지’의 시연. 벤치마크는 1회 실행(폭 우선)."),
        ("이념·가치 축이 없다", "과거 정책 갭 측정(시뮬 vs 당시 여론조사)에서 찬성률 "
         "MAE 27.4%p — 주범은 반대 진영 소실. 처방은 프롬프트가 아니라 데이터 강화."),
    ]
    y = 1.8
    for title, body in items:
        box(slide, 0.6, y, 12.13, 0.87, fill=CARD, line=LINE)
        tf = tb(slide, 0.88, y + 0.08, 11.6, 0.72)
        para(tf, [(title + "   ", {"bold": True, "color": NAVY, "size": 13}),
                  (body, {"size": 11.8, "color": INK})], first=True, after=0,
             lh=1.18)
        y += 0.97
    note(slide, "마무리 전에 한계를 저희 입으로 먼저 말씀드립니다. 미리랩은 미래를 "
         "맞히는 도구가 아니고, 그렇게 주장하지도 않습니다. 소득 판정은 휴리스틱 "
         "근사이고, '안 믿지만 받으면 좋지'를 분리할 실현 기대 측정축이 아직 없고, "
         "표본은 작습니다. 그리고 가장 큰 공백을 정직하게 — 과거 정책 4건을 당시 "
         "여론조사와 직접 맞대본 갭 측정에서 찬성률 평균 오차가 27.4%p였고, 주범은 "
         "이념·가치 축 부재로 인한 '반대 진영 소실'이었습니다. 이건 프롬프트 수정이 "
         "아니라 페르소나 데이터 강화로 풀 문제로 진단해 뒀습니다. 저희가 보여드린 건 "
         "'정확한 예측'이 아니라 — 페르소나에 근거하고, 상식 방향으로 구조화되고, "
         "기준을 먼저 등록하고 재는 '영향 시나리오'입니다. 이 정직한 선 긋기가 저희 "
         "설계의 일부입니다.")


def s19_closing(prs, sprites):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, 0, 0, SW, 0.18, fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
    t = tb(slide, 1.0, 1.15, SW - 2.0, 2.3)
    para(t, "“같은 정책, 다른 인생.”", size=40, bold=True, color=NAVY,
         first=True, after=10, align=PP_ALIGN.CENTER)
    para(t, [("미리랩은 그 갈림을 ", {"size": 18}),
             ("배포 전에", {"bold": True, "color": BLUE, "size": 18}),
             (" 보여줍니다", {"size": 18})], after=4, align=PP_ALIGN.CENTER)
    para(t, "— 누가 이해 못 하고 · 누가 신청 못 하고 · 어디서 막히는지",
         size=15, color=SUB, after=0, align=PP_ALIGN.CENTER)
    # 다음 단계
    nx = 0.7
    nexts = [
        ("이념·가치 축", "갭 27.4%p의 주범 ‘반대 진영 소실’ — 데이터 강화로 해소"),
        ("실현 기대 문항", "“안 믿지만 받으면 좋지”를 분리하는 설문 축 추가"),
        ("갭 리베이스라인", "Gemini 기준 재측정 — 개선 폭을 같은 자로 추적"),
        ("미리마을 확장", "누적 기억(B단계) · 적응 스케줄 · 페르소나 RAG"),
    ]
    for title, body in nexts:
        box(slide, nx, 3.75, 2.88, 1.35, fill=PALE_NAVY, line=NAVY, lw=0.75)
        tf = tb(slide, nx + 0.18, 3.9, 2.55, 1.1)
        para(tf, "다음 — " + title, size=12, bold=True, color=NAVY, first=True,
             after=3)
        para(tf, body, size=10.5, color=INK, after=0, lh=1.18)
        nx += 3.06
    # 스프라이트 + 감사
    keys = ["grandma", "minsu", "sua", "owner", "oldman"]
    x = (SW - len(keys) * 0.75) / 2
    for k in keys:
        p = sprites.get(k)
        if p:
            slide.shapes.add_picture(str(p), Inches(x), Inches(5.35),
                                     Inches(0.6), Inches(0.6))
            x += 0.75
    t2 = tb(slide, 1.0, 6.15, SW - 2.0, 0.9)
    para(t2, "감사합니다 — Q&A", size=22, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, first=True, after=4)
    para(t2, "데이터 nvidia/Nemotron-Personas-Korea (CC BY 4.0) · LangGraph · "
         "Streamlit · gemini-3-flash / gpt-4o-mini · Generative Agents (Park et al., 2023)",
         size=10, color=SUB, align=PP_ALIGN.CENTER, after=0)
    note(slide, "정리하면 — 미리랩은 실제 인구통계 기반 가상 시민에게 정책을 먼저 "
         "물어봐서, '같은 정책, 다른 인생'의 갈림을 배포 전에 보여주는 실험실입니다. "
         "다음 단계의 1순위는 이념·가치 축 — 과거 정책 갭 측정에서 확인된 '반대 진영 "
         "소실'(찬성률 MAE 27.4%p)을 페르소나 데이터 강화로 푸는 일입니다. 이어서 "
         "실현 기대 문항, Gemini 기준 갭 리베이스라인, 미리마을 확장이 준비돼 있습니다. "
         "들어주셔서 감사합니다. "
         "(Q&A 예상 질문: ①'실제 사람 반응과 맞는지 확인했나' → 정직하게: 과거 정책 "
         "4건을 당시 여론조사와 직접 맞대봤고 평균 오차 27.4%p, 주범(이념축 부재)과 "
         "처방(데이터 강화)까지 진단해 둔 상태라고 "
         "답변. Park et al.도 human eval과 비교했음을 인지하고 있다고 덧붙이면 강함 "
         "②비용은 → 시뮬 1회 1~2센트, 검증·벤치마크 전체 세트도 몇 달러 수준 "
         "③실제 입안자가 어떻게 쓰나 → 4장 컨셉의 사각지대 발견 + 정책 개선 탭의 "
         "종합 리포트(.md)로 답변)")


# ---------------------------------------------------------------- 메인
def main():
    sprites = prep_sprites()
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    s01_title(prs, sprites)
    s02_agenda(prs)
    s03_problem(prs)
    s04_concept(prs)
    s05_system(prs)
    s05b_dataset(prs)
    s06_persona(prs)
    s07_axis_a(prs)
    s08_axis_b(prs)
    s08b_axis3(prs)
    s09_scope(prs)
    s10_two_judges(prs)
    s11_billion(prs)
    s12_guards(prs)
    s13_validation(prs)
    s14_prereg(prs)
    s14b_bench(prs)
    s15_board(prs)
    s16_village1(prs, sprites)
    s17_village2(prs)
    s18_discipline(prs)
    s19_limits(prs)
    s19_closing(prs, sprites)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print("[OK] saved:", OUT)
    print("[OK] slides:", len(prs.slides.__iter__.__self__._sldIdLst))


if __name__ == "__main__":
    main()
